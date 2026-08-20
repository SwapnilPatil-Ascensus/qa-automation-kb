"""Shared design tokens and helpers for leadership deliverables."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1] / "2026-08-am-squad-leadership-update"
METRICS = ROOT / "data" / "leadership-metrics.json"
CHARTS = ROOT / "assets" / "charts"
DELIVERABLES = ROOT / "deliverables"

SUBTITLE = "April – August 2026"
TITLE = "QA Automation — AM Squad Leadership Update"

# Modern BI palette (PPTX)
INK = RGBColor(0x0A, 0x16, 0x28)
SURFACE = RGBColor(0xF0, 0xF4, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xA3, 0xE0)
TEAL = RGBColor(0x00, 0xB3, 0x88)
CORAL = RGBColor(0xFF, 0x6B, 0x35)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
NAVY = RGBColor(0x00, 0x32, 0x41)
TEXT = RGBColor(0x1A, 0x23, 0x32)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INSIGHT = RGBColor(0xE8, 0xEE, 0xF6)

# DOCX hex equivalents
DOCX_INK = "0A1628"
DOCX_CYAN = "00A3E0"
DOCX_TEAL = "00B388"
DOCX_CORAL = "FF6B35"
DOCX_VIOLET = "7C3AED"
DOCX_NAVY = "003241"
DOCX_INSIGHT = "E8EEF6"
DOCX_SURFACE = "F0F4F8"
DOCX_MUTED = "64748B"
DOCX_TEXT = "1A2332"

SW = Inches(10)
SH = Inches(7.5)
HEADER_H = Inches(0.72)
FOOTER_H = Inches(0.38)
FOOTER_H_NOTE = Inches(0.62)

REFS = [
    ("Jira — QA Automation project", "https://ascensuscollegesavings.atlassian.net/jira/software/projects/QA"),
    ("GitLab — api-test-automation", "https://gitlab.com/ascensus-gs/products/depot/qa-automation/api-test-automation"),
    ("GitLab — unite-test-automation (V2)", "https://gitlab.com/ascensus-gs/products/depot/qa-automation/automation"),
    ("GitLab — prime-test-automation (V3)", "https://gitlab.com/ascensus-gs/products/depot/qa-automation/prime-test-automation"),
    ("qTest — Automation Unite", "https://ascensus.qtestnet.com"),
    ("Jenkins — QA performance & nightly regression", "jenkinsqant1 (internal QA Jenkins)"),
]


def load_metrics() -> dict:
    with METRICS.open(encoding="utf-8") as f:
        return json.load(f)


def _rect(slide, left, top, width, height, fill, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    return shape


def _round_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def slide_surface(slide) -> None:
    _rect(slide, 0, 0, SW, SH, SURFACE)


def add_header(slide, title: str, tag: str = "") -> None:
    _rect(slide, 0, 0, SW, HEADER_H, INK)
    _rect(slide, 0, HEADER_H, SW, Inches(0.04), CYAN)
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.14), Inches(7.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if tag:
        badge = _round_rect(slide, Inches(8.15), Inches(0.18), Inches(1.55), Inches(0.38), CYAN)
        badge.text_frame.paragraphs[0].text = tag
        badge.text_frame.paragraphs[0].font.size = Pt(9)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide, source: str, num: int, footnote: str = "") -> None:
    fh = FOOTER_H_NOTE if footnote else FOOTER_H
    y = SH - fh
    _rect(slide, 0, y, SW, fh, CARD)
    _rect(slide, 0, y, SW, Inches(0.02), INSIGHT)
    row_y = y + Inches(0.07)
    if footnote:
        fb = slide.shapes.add_textbox(Inches(0.42), row_y, Inches(8.8), Inches(0.28))
        fp = fb.text_frame.paragraphs[0]
        fp.text = f"Note: {footnote}"
        fp.font.size = Pt(7.5)
        fp.font.italic = True
        fp.font.color.rgb = MUTED
        row_y += Inches(0.26)
    tb = slide.shapes.add_textbox(Inches(0.42), row_y, Inches(8.5), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.text = source
    p.font.size = Pt(7.5)
    p.font.color.rgb = MUTED
    nb = slide.shapes.add_textbox(Inches(9.2), row_y, Inches(0.6), Inches(0.22))
    np = nb.text_frame.paragraphs[0]
    np.text = str(num)
    np.font.size = Pt(8)
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.RIGHT


def add_kpi_card(slide, x, y, w, h, label, value, accent, sub: str = "") -> None:
    _round_rect(slide, x, y, w, h, CARD)
    _rect(slide, x, y, w, Inches(0.06), accent)
    lb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.14), w - Inches(0.2), Inches(0.35))
    lp = lb.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(9)
    lp.font.color.rgb = MUTED
    vb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.42), w - Inches(0.2), Inches(0.55))
    vp = vb.text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(22)
    vp.font.bold = True
    vp.font.color.rgb = TEXT
    if sub:
        sb = slide.shapes.add_textbox(x + Inches(0.12), y + h - Inches(0.32), w - Inches(0.2), Inches(0.28))
        sp = sb.text_frame.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(8)
        sp.font.color.rgb = MUTED


def add_insight_panel(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
    panel = _round_rect(slide, x, y, w, h, INSIGHT)
    _rect(slide, x, y, Inches(0.06), h, TEAL)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.35))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(11)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.48), w - Inches(0.3), h - Inches(0.55))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)


def slide_chart_insight(
    prs: Presentation,
    num: int,
    title: str,
    chart: str,
    insight_title: str,
    bullets: list[str],
    source: str,
    tag: str = "",
    chart_w: float = 5.35,
    footnote: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, title, tag)
    path = CHARTS / chart
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.35), Inches(0.95), width=Inches(chart_w))
    add_insight_panel(slide, Inches(5.85), Inches(0.95), Inches(3.8), Inches(5.55), insight_title, bullets)
    add_footer(slide, source, num, footnote)


def slide_hero(prs: Presentation, title: str, subtitle: str, kicker: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, INK)
    _rect(slide, 0, Inches(5.8), SW, Inches(1.7), NAVY)
    _rect(slide, 0, Inches(2.1), Inches(0.12), Inches(2.8), CYAN)
    if kicker:
        kb = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(8), Inches(0.4))
        kp = kb.text_frame.paragraphs[0]
        kp.text = kicker.upper()
        kp.font.size = Pt(11)
        kp.font.bold = True
        kp.font.color.rgb = CYAN
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.0), Inches(8.8), Inches(1.4))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    sb = slide.shapes.add_textbox(Inches(0.55), Inches(3.5), Inches(8), Inches(0.6))
    sp = sb.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = TEAL
    chips = ["V2 UI", "V3 UP", "API/MSC", "Perf", "Pipeline", "Standards"]
    for i, chip in enumerate(chips):
        cx = Inches(0.55 + (i % 3) * 3.1)
        cy = Inches(4.35 + (i // 3) * 0.55)
        c = _round_rect(slide, cx, cy, Inches(2.85), Inches(0.42), RGBColor(0x14, 0x28, 0x40))
        c.text_frame.paragraphs[0].text = chip
        c.text_frame.paragraphs[0].font.size = Pt(10)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_kpi_dashboard(prs: Presentation, num: int, m: dict) -> None:
    sc = m["scorecard"]
    j = m["jira"]["totals"]
    dc = m.get("data_confidence", {})
    ui = m.get("ui_inventory_scope", {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, "Executive Scorecard", "Live metrics")
    cards = [
        (Inches(0.35), "GitLab merges", str(sc["gitlab_merges"]), CYAN, "3 repos · Apr–Aug"),
        (Inches(2.55), "Story points", f"{j['story_points']:.0f}", TEAL, "Sprints 26.04–26.12"),
        (Inches(4.75), "V2 Stage1 nightly", str(sc["v2_nightly_methods"]), NAVY, "Aug 4 · ~12 mo build"),
        (Inches(6.95), "V3 Stage1 nightly", str(sc["v3_nightly_methods"]), VIOLET, "GitLab · Aug 4"),
        (Inches(0.35), Inches(2.55), "Perf cases", str(sc["perf_test_cases"]), CORAL, "Labels × plans"),
        (Inches(2.55), Inches(2.55), "MSC M2", sc["msc_m2_endpoints"], TEAL, "Endpoints"),
        (Inches(4.75), Inches(2.55), "MSC M1", sc["msc_m1_core"], CYAN, "Core endpoints"),
        (Inches(6.95), Inches(2.55), "Period delivery", str(dc.get("period_delivery_estimate", "—")), CORAL, "Est. cases Apr–Aug"),
    ]
    for item in cards:
        if len(item) == 5:
            x, label, val, accent, sub = item
            y = Inches(1.05)
        else:
            x, y, label, val, accent, sub = item
        add_kpi_card(slide, x, y, Inches(2.05), Inches(1.25), label, val, accent, sub)
    nb = slide.shapes.add_textbox(Inches(0.35), Inches(3.95), Inches(9.3), Inches(0.55))
    np = nb.text_frame.paragraphs[0]
    np.text = (
        "V2/V3 = Stage1 nightly snapshots only (built since Q2 2025). "
        "Excludes smoke, Stage 2/5, integrations, +33 CSR Actions."
    )
    np.font.size = Pt(8.5)
    np.font.italic = True
    np.font.color.rgb = MUTED
    add_insight_panel(
        slide, Inches(0.35), Inches(4.5), Inches(9.3), Inches(2.05),
        "How to read these numbers",
        dc.get("leadership_talking_points", [])[:5],
    )
    footnote = dc.get("scorecard_footnote") or ui.get("scorecard_footnote", "")
    add_footer(slide, "Sources: GitLab MR export · Jira AMSQUAD · Jenkins/GitLab nightly logs", num, footnote)


def slide_section_modern(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, Inches(0.45), SH, CYAN)
    _rect(slide, Inches(0.45), 0, SW - Inches(0.45), SH, INK)
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(8.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(8), Inches(0.6))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = TEAL


def slide_bullets_modern(
    prs: Presentation, num: int, title: str, bullets: list[str], source: str, tag: str = "", note: str = ""
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, title, tag)
    add_insight_panel(slide, Inches(0.35), Inches(0.95), Inches(9.3), Inches(5.5), "Key points", bullets)
    if note:
        nb = slide.shapes.add_textbox(Inches(0.35), Inches(6.55), Inches(9.3), Inches(0.4))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(9)
        np.font.italic = True
        np.font.color.rgb = MUTED
    add_footer(slide, source, num)


def slide_close(prs: Presentation, line1: str, line2: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, INK)
    _rect(slide, 0, Inches(3.2), SW, Inches(0.05), CYAN)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(8.8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = line1
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sb = slide.shapes.add_textbox(Inches(0.55), Inches(3.6), Inches(8), Inches(0.5))
    sp = sb.text_frame.paragraphs[0]
    sp.text = line2
    sp.font.size = Pt(14)
    sp.font.color.rgb = TEAL


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs
