#!/usr/bin/env python3
"""Generate modern BI-dashboard style leadership PPTX decks (executive + detailed)."""

from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1] / "2026-08-am-squad-leadership-update"
METRICS = ROOT / "data" / "leadership-metrics.json"
CHARTS = ROOT / "assets" / "charts"
DELIVERABLES = ROOT / "deliverables"
EXEC_OUT = DELIVERABLES / "AM-Squad-Leadership-Executive-Glimpse-Aug2026.pptx"
DETAIL_OUT = DELIVERABLES / "AM-Squad-Leadership-Detailed-Modern-Aug2026.pptx"

SUBTITLE = "April – August 2026"

# Modern BI palette
INK = RGBColor(0x0A, 0x16, 0x28)
SURFACE = RGBColor(0xF0, 0xF4, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xA3, 0xE0)
TEAL = RGBColor(0x00, 0xB3, 0x88)
CORAL = RGBColor(0xFF, 0x6B, 0x35)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
NAVY = RGBColor(0x00, 0x32, 0x41)
TEXT = RGBColor(0x1A, 0x23, 0x32)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INSIGHT = RGBColor(0xE8, 0xEE, 0xF6)

SW = Inches(10)
SH = Inches(7.5)
HEADER_H = Inches(0.72)
FOOTER_H = Inches(0.38)
FOOTER_H_NOTE = Inches(0.62)


def load_metrics() -> dict:
    with METRICS.open(encoding="utf-8") as f:
        return json.load(f)


def _rect(slide, left, top, width, height, fill, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    return shape


def _round_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def slide_surface(slide) -> None:
    _rect(slide, 0, 0, SW, SH, SURFACE)


def add_header(slide, title: str, tag: str = "") -> None:
    _rect(slide, 0, 0, SW, HEADER_H, INK)
    _rect(slide, 0, HEADER_H, SW, Inches(0.04), CYAN)
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.14), Inches(7.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if tag:
        badge = _round_rect(slide, Inches(8.15), Inches(0.18), Inches(1.55), Inches(0.38), CYAN)
        badge.text_frame.paragraphs[0].text = tag
        badge.text_frame.paragraphs[0].font.size = Pt(9)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide, source: str, num: int, footnote: str = "") -> None:
    fh = FOOTER_H_NOTE if footnote else FOOTER_H
    y = SH - fh
    _rect(slide, 0, y, SW, fh, CARD)
    _rect(slide, 0, y, SW, Inches(0.02), INSIGHT)
    row_y = y + Inches(0.07)
    if footnote:
        fb = slide.shapes.add_textbox(Inches(0.42), row_y, Inches(8.8), Inches(0.28))
        fp = fb.text_frame.paragraphs[0]
        fp.text = f"Note: {footnote}"
        fp.font.size = Pt(7.5)
        fp.font.italic = True
        fp.font.color.rgb = MUTED
        row_y += Inches(0.26)
    tb = slide.shapes.add_textbox(Inches(0.42), row_y, Inches(8.5), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.text = source
    p.font.size = Pt(7.5)
    p.font.color.rgb = MUTED
    nb = slide.shapes.add_textbox(Inches(9.2), row_y, Inches(0.6), Inches(0.22))
    np = nb.text_frame.paragraphs[0]
    np.text = str(num)
    np.font.size = Pt(8)
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.RIGHT


def add_kpi_card(slide, x, y, w, h, label, value, accent, sub: str = "") -> None:
    card = _round_rect(slide, x, y, w, h, CARD)
    _rect(slide, x, y, w, Inches(0.06), accent)
    lb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.14), w - Inches(0.2), Inches(0.35))
    lp = lb.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(9)
    lp.font.color.rgb = MUTED
    vb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.42), w - Inches(0.2), Inches(0.55))
    vp = vb.text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(22)
    vp.font.bold = True
    vp.font.color.rgb = TEXT
    if sub:
        sb = slide.shapes.add_textbox(x + Inches(0.12), y + h - Inches(0.32), w - Inches(0.2), Inches(0.28))
        sp = sb.text_frame.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(8)
        sp.font.color.rgb = MUTED


def add_insight_panel(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
    panel = _round_rect(slide, x, y, w, h, INSIGHT)
    _rect(slide, x, y, Inches(0.06), h, TEAL)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.35))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(11)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.48), w - Inches(0.3), h - Inches(0.55))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)


def slide_chart_insight(
    prs: Presentation,
    num: int,
    title: str,
    chart: str,
    insight_title: str,
    bullets: list[str],
    source: str,
    tag: str = "",
    chart_w: float = 5.35,
    footnote: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, title, tag)
    path = CHARTS / chart
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.35), Inches(0.95), width=Inches(chart_w))
    add_insight_panel(slide, Inches(5.85), Inches(0.95), Inches(3.8), Inches(5.55), insight_title, bullets)
    add_footer(slide, source, num, footnote)


def slide_hero(prs: Presentation, title: str, subtitle: str, kicker: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, INK)
    _rect(slide, 0, Inches(5.8), SW, Inches(1.7), NAVY)
    _rect(slide, 0, Inches(2.1), Inches(0.12), Inches(2.8), CYAN)
    if kicker:
        kb = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(8), Inches(0.4))
        kp = kb.text_frame.paragraphs[0]
        kp.text = kicker.upper()
        kp.font.size = Pt(11)
        kp.font.bold = True
        kp.font.color.rgb = CYAN
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.0), Inches(8.8), Inches(1.4))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    sb = slide.shapes.add_textbox(Inches(0.55), Inches(3.5), Inches(8), Inches(0.6))
    sp = sb.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = TEAL
    chips = ["V2 UI", "V3 UP", "API/MSC", "Perf", "Pipeline", "Standards"]
    for i, chip in enumerate(chips):
        cx = Inches(0.55 + (i % 3) * 3.1)
        cy = Inches(4.35 + (i // 3) * 0.55)
        c = _round_rect(slide, cx, cy, Inches(2.85), Inches(0.42), RGBColor(0x14, 0x28, 0x40))
        c.text_frame.paragraphs[0].text = chip
        c.text_frame.paragraphs[0].font.size = Pt(10)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_kpi_dashboard(prs: Presentation, num: int, m: dict) -> None:
    sc = m["scorecard"]
    j = m["jira"]["totals"]
    dc = m.get("data_confidence", {})
    ui = m.get("ui_inventory_scope", {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, "Executive Scorecard", "Live metrics")
    cards = [
        (Inches(0.35), "GitLab merges", str(sc["gitlab_merges"]), CYAN, "3 repos · Apr–Aug"),
        (Inches(2.55), "Story points", f"{j['story_points']:.0f}", TEAL, "Sprints 26.04–26.12"),
        (Inches(4.75), "V2 Stage1 nightly", str(sc["v2_nightly_methods"]), NAVY, "Aug 4 · ~12 mo build"),
        (Inches(6.95), "V3 Stage1 nightly", str(sc["v3_nightly_methods"]), VIOLET, "GitLab · Aug 4"),
        (Inches(0.35), Inches(2.55), "Perf cases", str(sc["perf_test_cases"]), CORAL, "Labels × plans"),
        (Inches(2.55), Inches(2.55), "MSC M2", sc["msc_m2_endpoints"], TEAL, "Endpoints"),
        (Inches(4.75), Inches(2.55), "MSC M1", sc["msc_m1_core"], CYAN, "Core endpoints"),
        (Inches(6.95), Inches(2.55), "Period delivery", str(dc.get("period_delivery_estimate", "—")), CORAL, "Est. cases Apr–Aug"),
    ]
    for item in cards:
        if len(item) == 5:
            x, label, val, accent, sub = item
            y = Inches(1.05)
        else:
            x, y, label, val, accent, sub = item
        add_kpi_card(slide, x, y, Inches(2.05), Inches(1.25), label, val, accent, sub)
    note = (
        "V2/V3 = Stage1 nightly snapshots only (built since Q2 2025). "
        "Excludes smoke, Stage 2/5, integrations, +33 CSR Actions."
    )
    nb = slide.shapes.add_textbox(Inches(0.35), Inches(3.95), Inches(9.3), Inches(0.55))
    np = nb.text_frame.paragraphs[0]
    np.text = note
    np.font.size = Pt(8.5)
    np.font.italic = True
    np.font.color.rgb = MUTED
    add_insight_panel(
        slide, Inches(0.35), Inches(4.5), Inches(9.3), Inches(2.05),
        "How to read these numbers",
        dc.get("leadership_talking_points", [])[:5],
    )
    footnote = dc.get("scorecard_footnote") or ui.get("scorecard_footnote", "")
    add_footer(
        slide,
        "Sources: GitLab MR export · Jira AMSQUAD · Jenkins/GitLab nightly logs",
        num,
        footnote or "Do not add period delivery (~1,212) to nightly inventory — different metrics.",
    )


def slide_section_modern(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, Inches(0.45), SH, CYAN)
    _rect(slide, Inches(0.45), 0, SW - Inches(0.45), SH, INK)
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(8.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(8), Inches(0.6))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = TEAL


def slide_bullets_modern(
    prs: Presentation,
    num: int,
    title: str,
    bullets: list[str],
    source: str,
    tag: str = "",
    footnote: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_surface(slide)
    add_header(slide, title, tag)
    add_insight_panel(slide, Inches(0.35), Inches(0.95), Inches(9.3), Inches(5.35), "Key points", bullets)
    add_footer(slide, source, num, footnote)


def build_executive(m: dict) -> None:
    sc = m["scorecard"]
    j = m["jira"]["totals"]
    dc = m.get("data_confidence", {})
    perf_base = sum(a["labels"] for a in m["perf_inventory"]["areas"])
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    n = 1

    slide_hero(prs, "QA Automation — AM Squad", f"Leadership Glimpse · {SUBTITLE}", "Executive dashboard")

    slide_kpi_dashboard(prs, n := n + 1, m)

    slide_chart_insight(
        prs, n := n + 1,
        "GitLab Delivery Velocity",
        "01-gitlab-mrs-by-month.png",
        "What this shows",
        [
            f"{sc['gitlab_merges']} merges to main across automation, prime-test-automation, and api-test-automation",
            "July peak: 38 merges during Unite MSC API sprint",
            "Code delivery metric — separate from test case counts",
        ],
        "GitLab MR export · Apr 1 – Aug 4, 2026",
        "GitLab",
        footnote="July = 38 merges (highest month). Counts merged code changes, not test cases or story points.",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "Monthly Automation Delivery",
        "08-monthly-automation-test-cases-added.png",
        "How to interpret",
        [
            f"~{dc.get('period_delivery_estimate', 1212)} est. new test cases delivered Apr–Aug (not cumulative)",
            "225 Jira stories closed × channel-specific averages",
            "Includes multi-plan, multi-env, pos/neg permutations",
            "Pre-Apr foundation sits in nightly inventory — chart shows when work landed",
        ],
        "Jira AMSQUAD Sprints 26.04–26.12 · period delivery estimate",
        "Automation",
        footnote="Period velocity only — do not add ~1,212 to scorecard inventory (592+442+323).",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "Jira Sprint Delivery",
        "09-jira-story-points-by-sprint.png",
        "Sprint outcomes",
        [
            f"{j['work_items_in_sprints']} work items · {j['story_points']:.0f} story points",
            f"{j['automation_bugs_logged']} automation-discovered defects logged",
            "Sustained delivery across 9 sprints in reporting window",
        ],
        "Jira AMSQUAD export · Sprints 26.04–26.12",
        "Jira",
        footnote="Story points = committed sprint work closed. Bugs = defects found via nightly regression triage.",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "V2 Legacy UI — Stage1 Nightly Snapshot",
        "04-v2-regression-by-module.png",
        "What this counts",
        [
            f"{sc['v2_nightly_methods']} test methods — Jenkins Stage1 nightly only (Aug 4)",
            "Built since Q2 2025 — not delivered in Apr–Aug alone",
            "Bars = methods per module; excludes smoke, Stage 2/5, CSR Actions (+33)",
            "CSR maintenance modules added Apr–Jul on top of earlier foundation",
        ],
        "Jenkins STAGE1-Daily-Unite-Prime-Regression · Aug 4 snapshot",
        "V2",
        footnote=(
            "Additional V2 coverage exists outside this snapshot: Stage 5 smoke, Stage 2 smoke, "
            "+33 CSR Actions (pending nightly wire), on-demand fast smoke."
        ),
    )

    slide_chart_insight(
        prs, n := n + 1,
        "V3 Universal Platform — Stage1 Nightly Snapshot",
        "11-v3-regression-by-module.png",
        "What this counts",
        [
            f"{sc['v3_nightly_methods']} test methods — GitLab Stage1 nightly only (Aug 4)",
            "Built since Q2 2025 — framework, CI/CD, suites accumulated over ~1 year",
            "UE 303 = scenarios × plan/traunch — not 303 unique scripts",
            "Entity suites expanding — not all in Aug 4 nightly log",
        ],
        "GitLab scheduled_regression_job · Aug 4 snapshot",
        "V3",
        footnote=(
            "Additional V3 coverage exists outside this snapshot: Stage 5 smoke (UE + IDP), "
            "integration XML profiles, Entity track — not added to 442."
        ),
    )

    slide_chart_insight(
        prs, n := n + 1,
        "API / Unite MSC — Coverage",
        "05-unite-msc-coverage.png",
        "MSC rescue",
        [
            f"M2 {sc['msc_m2_endpoints']} endpoints · M1 {sc['msc_m1_core']} core categories",
            "Delivered in ~50% of original ETA",
            "AI-assisted migration + canonical TestNG framework",
            "P0: GitLab nightly scheduling (QA-1405)",
        ],
        "api-test-automation repo inventory",
        "API",
        footnote="M2 = full mobile endpoint catalog. M1 ~86% of 29 core categories — master suite still in progress.",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "Performance — Test Case Inventory",
        "12-perf-test-case-inventory.png",
        "Perf counting model",
        [
            f"{perf_base} base transaction flows → {sc['perf_test_cases']} plan-expanded cases",
            "Bars = flows; line = cases after × plan matrix",
            "IDP: 15 labels × 7 plans = 105 cases alone",
            "Barcode SYN-443 delivered in ~1 week",
        ],
        "Perf inventory model · Jenkins + BlazeMeter",
        "Perf",
        footnote="4 Jenkins scenarios schedule runs — not 323 separate jobs. Plan expansion is standard perf counting.",
    )

    slide_bullets_modern(
        prs, n := n + 1,
        "Leadership Asks & Next Steps",
        [
            "Roadmap visibility — engage AM Squad at SDLC start, not sign-off deadline",
            "Administrative capacity — free lead for architecture & AI tooling",
            "P0: MSC GitLab nightly · M1 master suite · enrollment API",
            f"~{sc['release_automation_pct']}% release validation automated (17 FTE → 2 FTE equivalent)",
        ],
        "Full detail: AM-Squad-Leadership-Detailed-Modern-Aug2026.pptx",
        "Asks",
        footnote="Team formed Q2 2025; Apr–Aug is the 5-month delivery window on top of ~12 months of build.",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, INK)
    _rect(slide, 0, Inches(3.2), SW, Inches(0.05), CYAN)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(8.8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "QA Automation AM Squad"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sb = slide.shapes.add_textbox(Inches(0.55), Inches(3.6), Inches(8), Inches(0.5))
    sp = sb.text_frame.paragraphs[0]
    sp.text = "Delivering across six automation channels · Team since Q2 2025"
    sp.font.size = Pt(14)
    sp.font.color.rgb = TEAL

    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    prs.save(EXEC_OUT)
    print(f"Wrote {EXEC_OUT}")


def build_detailed(m: dict) -> None:
    sc = m["scorecard"]
    j = m["jira"]["totals"]
    dc = m.get("data_confidence", {})
    perf_base = sum(a["labels"] for a in m["perf_inventory"]["areas"])
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    n = 1

    slide_hero(prs, "QA Automation — AM Squad", f"Detailed Leadership Update · {SUBTITLE}", "Full portfolio view")
    slide_kpi_dashboard(prs, n := n + 1, m)

    slide_bullets_modern(
        prs, n := n + 1,
        "Executive Headline",
        [
            f"{sc['gitlab_merges']} merged changes to main across 3 automation repositories",
            f"{j['story_points']:.0f} Jira story points · {j['work_items_in_sprints']} work items · {j['automation_bugs_logged']} bugs found",
            f"~{sc['release_automation_pct']}% of monthly release validations automated (17 FTE → 2 FTE)",
            "Unite MSC rescued — Mobile 2 at 25/25 endpoints, ~50% ETA savings",
            "Six parallel tracks: V2 · V3 · API/MSC · Perf · Pipeline · Standards",
        ],
        "Executive summary",
        "Overview",
        footnote="Apr–Aug reporting window on a team formed Q2 2025 — acceleration phase, not greenfield.",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "GitLab Delivery Velocity",
        "01-gitlab-mrs-by-month.png",
        "What this shows",
        [f"{sc['gitlab_merges']} total merges", "July peak: 38 merges during MSC sprint", "Separate from test case counts"],
        "GitLab MR export · Apr 1 – Aug 4, 2026",
        "GitLab",
        footnote="July = 38 merges (highest month). Merge count ≠ test cases or story points.",
    )
    slide_chart_insight(
        prs, n := n + 1,
        "Monthly Automation Delivery",
        "08-monthly-automation-test-cases-added.png",
        "Period delivery (not inventory)",
        dc.get("leadership_talking_points", [])[:4],
        "Jira AMSQUAD · period delivery estimate",
        "Automation",
        footnote="Estimated from 225 closed Jira stories — period velocity only, not added to nightly inventory.",
    )
    slide_chart_insight(
        prs, n := n + 1,
        "Jira Sprint Delivery",
        "09-jira-story-points-by-sprint.png",
        "Sprint outcomes",
        [f"{j['work_items_in_sprints']} items · {j['story_points']:.0f} SP", f"{j['automation_bugs_logged']} bugs logged"],
        "Jira AMSQUAD Sprints 26.04–26.12",
        "Jira",
        footnote="Story points = committed sprint work closed across 9 sprints in the reporting window.",
    )
    slide_chart_insight(
        prs, n := n + 1,
        "Automation Defects Discovered",
        "10-jira-automation-bugs-by-sprint.png",
        "Quality signal",
        ["Defects found via nightly regression triage", "Fed into automation bug lifecycle standard"],
        "Jira AMSQUAD",
        "Quality",
        footnote="Bugs logged by automation team from nightly failures — not all production defects.",
    )

    slide_bullets_modern(
        prs, n := n + 1,
        "Beyond the Metrics",
        [
            "Framework architecture & canonical repo design",
            "AI-accelerated docs, Postman, TestNG boilerplate",
            "qTest master suite & automation bug lifecycle standard",
            "Pipeline/DevOps co-design & module switches",
            "Cross-team emergency support & release triage",
        ],
        "Additional value not captured in merge or test-count metrics",
        "Value",
        footnote="These investments are not fully captured in delivery charts or merge counts.",
    )

    slide_section_modern(prs, "V2 Legacy UI Automation", "Jenkins Stage1 nightly · built since Q2 2025")
    slide_chart_insight(
        prs, n := n + 1,
        "V2 Stage1 Nightly Snapshot",
        "04-v2-regression-by-module.png",
        "What this counts",
        [
            f"{sc['v2_nightly_methods']} methods — Stage1 primary nightly only (Aug 4)",
            "Built since Q2 2025 — not Apr–Aug alone",
            "Excludes smoke, Stage 2/5, +33 CSR Actions pending wire",
            "CSR maintenance added Apr–Jul on earlier foundation",
        ],
        "Jenkins STAGE1-Daily-Unite-Prime-Regression · Aug 4",
        "V2",
        footnote="Additional V2 coverage: Stage 5 smoke, Stage 2 smoke, on-demand fast smoke, +33 CSR Actions.",
    )
    slide_bullets_modern(
        prs, n := n + 1,
        "V2 Business Value",
        [
            "Core member journeys covered nightly before release",
            "CSR maintenance gaps closed — high-risk flows automated",
            "Enrollments/login triage separates env vs product defects",
        ],
        "V2 value",
        "V2",
    )

    slide_section_modern(prs, "V3 Universal Platform", "GitLab CI nightly · built since Q2 2025")
    slide_chart_insight(
        prs, n := n + 1,
        "V3 Stage1 Nightly Snapshot",
        "11-v3-regression-by-module.png",
        "What this counts",
        [
            f"{sc['v3_nightly_methods']} methods — GitLab Stage1 nightly only (Aug 4)",
            "Framework + CI/CD accumulated over ~1 year",
            "UE 303 = scenarios × plan/traunch matrix",
            "Entity suites expanding — not all in Aug 4 log",
        ],
        "GitLab scheduled_regression_job · Aug 4",
        "V3",
        footnote="Additional V3 coverage: Stage 5 smoke (UE + IDP), integration profiles, Entity track — not in 442.",
    )
    slide_bullets_modern(
        prs, n := n + 1,
        "V3 Delivery Highlights",
        [
            "Entity registration/login suites expanded",
            "IDP open-account and member withdrawal regression",
            "Flaky-test stabilization and web registration flows",
            "Stage 5 smoke suites for UE + IDP",
        ],
        "V3 highlights",
        "V3",
    )

    slide_section_modern(prs, "API / Unite MSC", "Rescued · accelerated")
    slide_chart_insight(
        prs, n := n + 1,
        "MSC Endpoint Coverage",
        "05-unite-msc-coverage.png",
        "MSC status",
        [f"M2 {sc['msc_m2_endpoints']}", f"M1 {sc['msc_m1_core']}", "~50% ETA savings", "P0: GitLab nightly (QA-1405)"],
        "api-test-automation",
        "MSC",
        footnote="M2 = full mobile endpoint catalog. M1 ~86% of 29 core categories — master suite in progress.",
    )
    slide_chart_insight(
        prs, n := n + 1,
        "API Module Breakdown",
        "13-api-regression-by-module.png",
        "M1 categories",
        ["Auth, profile, biometric, device, bank", "Master suite in progress", "Endpoints × branding plans"],
        "api-test-automation",
        "API",
        footnote="OKD non-IDP + NYD/NMD IDP brandings expand endpoint permutations beyond raw count.",
    )

    slide_section_modern(prs, "Performance Testing", "Labels × plans model")
    slide_chart_insight(
        prs, n := n + 1,
        "Perf Test Case Inventory",
        "12-perf-test-case-inventory.png",
        "Why the numbers look large",
        [
            f"{perf_base} base transaction flows → {sc['perf_test_cases']} plan-expanded cases",
            "Bars = flows; line = cases after × plan matrix",
            "4 Jenkins scenarios schedule these permutations",
            "IDP alone: 15 labels × 7 plans = 105 cases",
        ],
        "Perf inventory",
        "Perf",
        footnote="4 Jenkins jobs schedule runs — not 323 separate jobs. Barcode SYN-443 delivered in ~1 week.",
    )
    slide_bullets_modern(
        prs, n := n + 1,
        "Performance Value",
        [
            "Repeatable baselines — no manual re-run each release",
            "Evidence for platform team pre/post patch comparison",
            "Department perf DoD and BlazeMeter reporting standard",
        ],
        "Perf value",
        "Perf",
    )

    slide_chart_insight(
        prs, n := n + 1,
        "Investment Allocation",
        "06-work-allocation-index.png",
        "Effort mix",
        ["MSC API 35%", "V2 UI 20%", "V3 UP 15%", "Perf 12%", "Pipeline/Standards remainder"],
        "Squad estimate Apr–Jul",
        "Portfolio",
        footnote="Relative effort index — not headcount FTE. MSC sprint drove Jun–Jul concentration.",
    )
    slide_chart_insight(
        prs, n := n + 1,
        "Release Automation Impact",
        "07-release-automation-impact.png",
        "Business value",
        [f"~{sc['release_automation_pct']}% automated", "17 FTE → 2 FTE equivalent", "Team focuses on triage + net-new"],
        "Release validation model",
        "Value",
        footnote="FTE equivalent = manual validation hours replaced by nightly regression, not headcount reduction.",
    )

    slide_bullets_modern(
        prs, n := n + 1,
        "AI Acceleration",
        [
            "MSC migration agents — Postman, docs, TestNG boilerplate",
            "Automation bug lifecycle — Cursor skill + standardized reporting",
            "Chart & metrics generators — reproducible leadership packs",
            "Coverage intelligence foundation for monthly dashboard",
        ],
        "AI tooling portfolio",
        "AI",
    )
    slide_bullets_modern(
        prs, n := n + 1,
        "Roadmap Q3–Q4 2026",
        [
            "P0: MSC GitLab nightly · M1 master suite",
            "P1: MSC enrollment API · CSR Actions expansion (+33 scenarios)",
            "P2: Entity V3 nightly · automated dashboard",
            "Engage AM Squad at SDLC start, not sign-off deadline",
        ],
        "Roadmap",
        "Roadmap",
    )
    slide_bullets_modern(
        prs, n := n + 1,
        "Leadership Asks",
        [
            "Roadmap visibility at SDLC start",
            "Admin capacity for architecture lead & AI tooling",
            "30-min live walkthrough recommended for Q&A",
        ],
        "Asks",
        "Asks",
        footnote="Concise executive view: AM-Squad-Leadership-Executive-Glimpse-Aug2026.pptx",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, INK)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(3.0), Inches(8.8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    prs.save(DETAIL_OUT)
    print(f"Wrote {DETAIL_OUT}")


def main() -> None:
    m = load_metrics()
    build_executive(m)
    build_detailed(m)


if __name__ == "__main__":
    main()
