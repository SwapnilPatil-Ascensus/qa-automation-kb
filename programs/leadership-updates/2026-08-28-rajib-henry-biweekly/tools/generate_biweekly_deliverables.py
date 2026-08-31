#!/usr/bin/env python3
"""Generate Rajib/Henry biweekly leadership status DOCX + premium UX PPTX."""

from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data" / "biweekly-metrics.json"
ASSETS = ROOT / "assets" / "charts"
DELIVERABLES = ROOT / "deliverables"
M1_CSV = (
    Path(__file__).resolve().parents[3]
    / "unite-msc"
    / "api-test-automation"
    / "mappings"
    / "mobile1-endpoint-current-state.csv"
)  # reference: 26 in-scope business endpoints

# Ascensus-inspired executive palette
INK = PptRGB(0x0A, 0x16, 0x28)
NAVY = PptRGB(0x00, 0x32, 0x41)
CYAN = PptRGB(0x00, 0xA3, 0xE0)
TEAL = PptRGB(0x00, 0xB3, 0x88)
CORAL = PptRGB(0xFF, 0x6B, 0x35)
VIOLET = PptRGB(0x7C, 0x3A, 0xED)
WHITE = PptRGB(0xFF, 0xFF, 0xFF)
SURFACE = PptRGB(0xF0, 0xF4, 0xF8)
MUTED = PptRGB(0x64, 0x74, 0x8B)
TEXT = PptRGB(0x1A, 0x23, 0x32)
GREEN = PptRGB(0x16, 0xA3, 0x4A)
AMBER = PptRGB(0xD9, 0x77, 0x06)

SW = PptInches(13.333)
SH = PptInches(7.5)
MARGIN = PptInches(0.6)

DOCX_NAVY = RGBColor(0x00, 0x32, 0x41)
DOCX_TEAL = RGBColor(0x00, 0xB3, 0x88)
DOCX_GRAY = RGBColor(0x64, 0x74, 0x8B)


def load_data() -> dict:
    with DATA.open(encoding="utf-8") as f:
        return json.load(f)


def save_chart(fig, name: str) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / name
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_msc_progress(m: dict) -> Path:
    u = m["unite_msc"]
    labels = ["Mobile 2", "Mobile 1", "Enrollment"]
    values = [u["mobile2"]["pct"], u["mobile1"]["pct"], u["enrollment"]["pct"]]
    colors = [TEAL, CYAN, CORAL]
    fig, ax = plt.subplots(figsize=(8, 4.5))
    y = np.arange(len(labels))
    bars = ax.barh(y, values, color=[f"#{c:06x}" for c in [0x00B388, 0x00A3E0, 0xFF6B35]], height=0.55)
    ax.set_xlim(0, 110)
    ax.set_yticks(y, labels)
    ax.set_xlabel("Completion %")
    ax.set_title("Unite MSC — Program Progress (Sprint 26.14)", fontweight="bold", color="#003241")
    for bar, val in zip(bars, values):
        ax.text(val + 2, bar.get_y() + bar.get_height() / 2, f"{val}%", va="center", fontweight="bold")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", alpha=0.2)
    return save_chart(fig, "01-msc-progress.png")


def chart_sprint_mix(m: dict) -> Path:
    mix = m["sprint_focus_mix"]
    labels = mix["labels"]
    x = np.arange(len(labels))
    w = 0.25
    fig, ax = plt.subplots(figsize=(9, 4.8))
    ax.bar(x - w, mix["sprint_26_13"], w, label="Sprint 26.13", color="#003241")
    ax.bar(x, mix["sprint_26_14"], w, label="Sprint 26.14 (current)", color="#00A3E0")
    ax.bar(x + w, mix["sprint_26_15_plan"], w, label="Sprint 26.15 (plan)", color="#00B388")
    ax.set_xticks(x, labels, rotation=15, ha="right")
    ax.set_ylabel("Relative effort %")
    ax.set_title("Sprint Focus — 26.13 · 26.14 · 26.15 (plan)", fontweight="bold")
    ax.legend(fontsize=8, loc="upper right")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", alpha=0.2)
    return save_chart(fig, "02-sprint-focus-mix.png")


def chart_msc_endpoint_coverage(m: dict) -> Path:
    u = m["unite_msc"]
    labels = ["Mobile 2", "Mobile 1"]
    totals = [u["mobile2"]["total"], u["mobile1"]["total"]]
    done = [u["mobile2"]["automated"], u["mobile1"]["automated"]]
    fig, ax = plt.subplots(figsize=(8, 4.2))
    y = np.arange(len(labels))
    bars = ax.barh(y, [100, 100], color=["#00B388", "#00A3E0"], height=0.5)
    ax.set_xlim(0, 110)
    ax.set_yticks(y, labels)
    ax.set_xlabel("In-scope endpoints automated (%)")
    ax.set_title("Unite MSC — Endpoint Coverage (in-scope business endpoints)", fontweight="bold", color="#003241")
    for bar, d, t in zip(bars, done, totals):
        ax.text(102, bar.get_y() + bar.get_height() / 2, f"{d}/{t}", va="center", fontweight="bold", fontsize=10)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", alpha=0.2)
    fig.text(0.02, 0.02, m["unite_msc"]["footnote"], fontsize=7, color="#64748B", wrap=True)
    return save_chart(fig, "04-msc-endpoint-coverage.png")


def chart_enrollment_endpoints(m: dict) -> Path:
    e = m["enrollment_coverage"]
    labels = e["labels"]
    coded = e["coded"]
    remaining = e["remaining"]
    x = np.arange(len(labels))
    w = 0.35
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.bar(x - w / 2, coded, w, label="Coded", color="#00B388")
    ax.bar(x + w / 2, remaining, w, label="Remaining", color="#FF6B35")
    ax.set_xticks(x, labels, rotation=12, ha="right")
    ax.set_ylabel("Endpoints / work items")
    ax.set_title("Enrollment — Endpoint Coverage vs Excel Catalog", fontweight="bold", color="#003241")
    ax.legend(fontsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", alpha=0.2)
    return save_chart(fig, "05-enrollment-endpoint-coverage.png")


def chart_capacity_timeline(m: dict) -> Path:
    phases = ["Now\n(26.14)", "Coding done\n(26.14 end)", "Wrap-up\n(1 resource)", "Q4+"]
    msc = [92, 98, 100, 100]
    other = [30, 55, 75, 85]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.plot(phases, msc, marker="o", linewidth=2.5, color="#00B388", label="MSC program %")
    ax.plot(phases, other, marker="s", linewidth=2.5, color="#7C3AED", label="Capacity for new work %")
    ax.set_ylim(0, 110)
    ax.set_title("Capacity Shift — MSC Wrap-up → New Priorities", fontweight="bold")
    ax.legend(loc="lower right")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.2)
    return save_chart(fig, "03-capacity-timeline.png")


def generate_charts(m: dict) -> None:
    chart_msc_progress(m)
    chart_sprint_mix(m)
    chart_capacity_timeline(m)
    chart_msc_endpoint_coverage(m)
    chart_enrollment_endpoints(m)


# ── PPTX helpers ─────────────────────────────────────────────────────────────

def _rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _round(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def slide_bg(slide, dark: bool = False) -> None:
    _rect(slide, 0, 0, SW, SH, INK if dark else SURFACE)


def slide_header(slide, title: str, subtitle: str = "", tag: str = "") -> None:
    _rect(slide, 0, 0, SW, PptInches(1.15), INK)
    _rect(slide, 0, PptInches(1.15), SW, PptInches(0.05), CYAN)
    tb = slide.shapes.add_textbox(MARGIN, PptInches(0.28), PptInches(9.5), PptInches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(MARGIN, PptInches(0.72), PptInches(9), PptInches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = PptPt(12)
        sp.font.color.rgb = TEAL
    if tag:
        badge = _round(slide, PptInches(11.2), PptInches(0.35), PptInches(1.6), PptInches(0.42), CYAN)
        badge.text_frame.paragraphs[0].text = tag
        badge.text_frame.paragraphs[0].font.size = PptPt(9)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_footer(slide, note: str, num: int) -> None:
    y = PptInches(7.05)
    _rect(slide, 0, y, SW, PptInches(0.45), WHITE)
    tb = slide.shapes.add_textbox(MARGIN, y + PptInches(0.1), PptInches(11), PptInches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = note
    p.font.size = PptPt(8)
    p.font.color.rgb = MUTED
    nb = slide.shapes.add_textbox(PptInches(12.5), y + PptInches(0.1), PptInches(0.5), PptInches(0.3))
    np = nb.text_frame.paragraphs[0]
    np.text = str(num)
    np.font.size = PptPt(9)
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.RIGHT


def add_kpi_card(slide, x, y, w, h, label, value, accent, sub=""):
    card = _round(slide, x, y, w, h, WHITE)
    _rect(slide, x, y, w, PptInches(0.07), accent)
    lb = slide.shapes.add_textbox(x + PptInches(0.15), y + PptInches(0.18), w - PptInches(0.25), PptInches(0.35))
    lb.text_frame.paragraphs[0].text = label
    lb.text_frame.paragraphs[0].font.size = PptPt(10)
    lb.text_frame.paragraphs[0].font.color.rgb = MUTED
    vb = slide.shapes.add_textbox(x + PptInches(0.15), y + PptInches(0.5), w - PptInches(0.25), PptInches(0.55))
    vb.text_frame.paragraphs[0].text = value
    vb.text_frame.paragraphs[0].font.size = PptPt(26)
    vb.text_frame.paragraphs[0].font.bold = True
    vb.text_frame.paragraphs[0].font.color.rgb = TEXT
    if sub:
        sb = slide.shapes.add_textbox(x + PptInches(0.15), y + h - PptInches(0.38), w - PptInches(0.25), PptInches(0.3))
        sb.text_frame.paragraphs[0].text = sub
        sb.text_frame.paragraphs[0].font.size = PptPt(8)
        sb.text_frame.paragraphs[0].font.color.rgb = MUTED


def slide_bullets(slide, title: str, bullets: list[str], x, y, w, h, accent=TEAL, font_size=9.5):
    panel = _round(slide, x, y, w, h, PptRGB(0xE8, 0xEE, 0xF6))
    _rect(slide, x, y, PptInches(0.07), h, accent)
    tb = slide.shapes.add_textbox(x + PptInches(0.22), y + PptInches(0.12), w - PptInches(0.35), PptInches(0.35))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = PptPt(12)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = NAVY
    body = slide.shapes.add_textbox(x + PptInches(0.22), y + PptInches(0.48), w - PptInches(0.35), h - PptInches(0.55))
    tf = body.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = PptPt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = PptPt(3)
        p.line_spacing = 1.05


def slide_categorized_asks(slide, asks: dict, x, y, w, h):
    panel = _round(slide, x, y, w, h, PptRGB(0xE8, 0xEE, 0xF6))
    _rect(slide, x, y, PptInches(0.07), h, CORAL)
    col_w = (w - PptInches(0.5)) / 2
    col_h = (h - PptInches(0.35)) / 3
    items = list(asks.items())
    for idx, (category, bullets) in enumerate(items):
        col = idx % 2
        row = idx // 2
        cx = x + PptInches(0.22) + col * (col_w + PptInches(0.15))
        cy = y + PptInches(0.12) + row * (col_h + PptInches(0.08))
        tb = slide.shapes.add_textbox(cx, cy, col_w, col_h)
        tf = tb.text_frame
        tf.word_wrap = True
        hp = tf.paragraphs[0]
        hp.text = category
        hp.font.size = PptPt(10)
        hp.font.bold = True
        hp.font.color.rgb = NAVY
        hp.space_after = PptPt(2)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = PptPt(7.5)
            p.font.color.rgb = TEXT
            p.space_after = PptPt(1)
            p.line_spacing = 1.0


def slide_chart_page(prs, n, title, subtitle, chart, bullets, note, tag="", panel_h=PptInches(5.35)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, title, subtitle, tag)
    path = ASSETS / chart
    if path.exists():
        slide.shapes.add_picture(str(path), MARGIN, PptInches(1.45), width=PptInches(6.2))
    slide_bullets(slide, "Key messages", bullets, PptInches(7.0), PptInches(1.45), PptInches(5.7), panel_h, font_size=9)
    slide_footer(slide, note, n)


def build_pptx(m: dict) -> Path:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    n = 0
    u = m["unite_msc"]
    p = m["pulse"]
    ui = m["ui_automation"]
    perf = m["performance"]
    cap = m["capacity"]
    meeting = datetime.strptime(m["meeting_date"], "%Y-%m-%d").strftime("%B %d, %Y")

    # 1 Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, dark=True)
    _rect(slide, 0, PptInches(2.0), PptInches(0.14), PptInches(3.2), CYAN)
    t = slide.shapes.add_textbox(MARGIN, PptInches(2.1), PptInches(11), PptInches(1.2))
    t.text_frame.paragraphs[0].text = "QA Automation — AM Squad"
    t.text_frame.paragraphs[0].font.size = PptPt(40)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    s = slide.shapes.add_textbox(MARGIN, PptInches(3.35), PptInches(11), PptInches(0.6))
    s.text_frame.paragraphs[0].text = f"Biweekly Status · {meeting}"
    s.text_frame.paragraphs[0].font.size = PptPt(18)
    s.text_frame.paragraphs[0].font.color.rgb = TEAL
    a = slide.shapes.add_textbox(MARGIN, PptInches(4.2), PptInches(11), PptInches(0.5))
    a.text_frame.paragraphs[0].text = "Rajib · Henry · Persistent Delivery Managers"
    a.text_frame.paragraphs[0].font.size = PptPt(13)
    a.text_frame.paragraphs[0].font.color.rgb = MUTED
    n += 1

    # 2 Pulse
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Executive Pulse", p["sprint_scope"], "At a glance")
    cards = [
        (MARGIN, "GitLab merges", str(p["gitlab_merges"]), CYAN, p["gitlab_merges_note"]),
        (PptInches(3.35), "Sprint 26.14 items", str(p["jira_sprint_items_current"]), TEAL, p["jira_sprint_items_note"]),
        (PptInches(6.05), "MSC in-scope APIs", p["msc_in_scope_endpoints"], NAVY, p["msc_in_scope_note"]),
        (PptInches(8.75), "Regression bugs", str(p["automation_bugs_logged"]), CORAL, p["automation_bugs_note"]),
        (PptInches(11.45), "ETA track", p["eta_delivery_rate"], GREEN, "3-sprint delivery window"),
    ]
    for x, lb, val, ac, sub in cards:
        add_kpi_card(slide, x, PptInches(1.55), PptInches(2.55), PptInches(1.35), lb, val, ac, sub)
    slide_bullets(slide, "What leadership should know", [
        "Scope: Sprints 26.13 and 26.14 (current) — not Apr–Aug historical totals",
        "Mobile 1 and Mobile 2: 100% — sign-off ready; enrollment core E2E ~95% (19/20)",
        "Coding complete ≠ program closed: 1 resource ~1 sprint for docs, qTest, Bruno, KT",
        "DB refresh completed this week — test data and SQL stabilized for enrollment runs",
    ], MARGIN, PptInches(3.2), PptInches(12.1), PptInches(3.6), font_size=9.5)
    slide_footer(slide, "Sources: Jira AMSQUAD board · GitLab (3 repos) · Jenkins/GitLab nightly logs", n := n + 1)

    # 3 Sprint focus
    slide_chart_page(prs, n := n + 1, "Sprint Delivery Focus",
        "Sprints 26.13 · 26.14 (current) · 26.15 plan", "02-sprint-focus-mix.png",
        [
            "Sprint 26.13: Enrollment API slices, M1/M2 sign-off prep, GitLab nightly delivered",
            "Sprint 26.14: Enrollment wizard completion, DB refresh, submit + Stage1 E2E",
            "Sprint 26.15 (plan): MSC wrap-up (1 resource) + Atlas intake / V2-V3 backlog",
            "Perf: Preeti on contribution JMX; M1+enrollment perf scripts next sprint",
        ], "Relative effort index — not FTE hours", "Delivery", panel_h=PptInches(5.5))

    # 4 MSC
    m1 = u["mobile1"]
    enr = u["enrollment"]
    slide_chart_page(prs, n := n + 1, "Unite MSC — Program Status",
        f"Mobile 2 100% · Mobile 1 100% · Enrollment {enr['pct']}%", "01-msc-progress.png",
        [
            f"Mobile 2: {u['mobile2']['automated']}/{u['mobile2']['total']} endpoints — sign-off ready",
            f"Mobile 1: {m1['automated']}/{m1['total']} in-scope operations — sign-off ready",
            f"Enrollment: {enr['core_e2e']} core E2E ({enr['pct']}%) — {enr['test_classes']} test classes in repo",
            f"Remaining coding: {', '.join(enr['remaining_coding'])}",
            f"Target: {u['target_completion']}",
        ], u["footnote"], "MSC", panel_h=PptInches(5.55))

    # 5 Enrollment coverage
    slide_chart_page(prs, n := n + 1, "Enrollment — Endpoint Coverage",
        f"Core E2E {enr['core_e2e']} · Catalog {enr['catalog']}", "05-enrollment-endpoint-coverage.png",
        [
            "17 TestNG classes: smoke GETs + full wizard through allocations + optional helpers",
            "Submit (review-confirm-entered) — target end Sprint 26.14 for account creation",
            "Deferred next sprint: subsequent enrollment, Vanguard, Upromise, OAuth (research)",
            enr["db_refresh"],
            "Postman/Excel source of truth synced in KB — Endpoint catalog + E2E collection",
        ], "Repo: api-test-automation/mobile/enrollment", "Enrollment", panel_h=PptInches(5.55))
    slide_chart_page(prs, n := n + 1, "Unite MSC — Endpoint Coverage",
        "M2 25/25 · M1 26/26 in-scope", "04-msc-endpoint-coverage.png",
        [
            "All in-scope business endpoints automated — 100% for Mobile 1 and Mobile 2",
            "Destructive endpoints run in smoke/integration — not duplicated in master regression",
            "Stage1: IDP + non-IDP paths validated; QC4 hybrid proof pending env stability",
            "GitLab nightly + ENVP pipeline integration (QA-1405, QA-1544–1549)",
            "Need leadership to name sign-off owners for M1, M2, pipeline, and perf",
        ], "Repo: programs/unite-msc/api-test-automation/mappings/", "MSC", panel_h=PptInches(5.55))

    # 6 Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Pipeline & Environments", "Hybrid QC4 + Stage1 model", "Infrastructure")
    slide_bullets(slide, "Current state", [
        u["pipeline"]["stage1"],
        u["pipeline"]["qc4"],
        u["pipeline"]["gitlab_nightly"],
        "QC4 has no non-IDP plants in lower env today — need product owner to confirm test plans",
        "Pipeline is not 'done' until QC4 proof passes — design is ready",
    ], MARGIN, PptInches(1.5), PptInches(5.8), PptInches(5.2), CYAN)
    slide_bullets(slide, "What we need", [
        "DevOps: GitLab nightly job for MSC regression (in progress)",
        "Leadership: identify non-IDP plan contacts for lower-environment testing",
        "Sign-off: names for M1/M2 acceptance before we schedule review sessions",
    ], PptInches(6.7), PptInches(1.5), PptInches(5.9), PptInches(5.2), CORAL)
    slide_footer(slide, "Aligned with Jul 2025 hybrid decision — Stage1 primary until QC4 stabilizes", n := n + 1)

    # 7 UI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "V2 & V3 UI Automation", ui["scope_note"], "UI")
    slide_bullets(slide, "V3 (GitLab) — maintenance + stabilization", ui["v3_highlights"],
                  MARGIN, PptInches(1.5), PptInches(5.9), PptInches(2.3), VIOLET, font_size=9)
    slide_bullets(slide, "V2 (Jenkins) — sunset discussion needed", ui["v2_highlights"] + [ui["migration_discussion"]],
                  PptInches(6.7), PptInches(1.5), PptInches(5.9), PptInches(2.3), NAVY, font_size=9)
    bug_bullets = [f"{b['key']}: {b['summary']} ({b['sprint']})" for b in m["recent_bugs"]]
    slide_bullets(slide, "Regression bugs — Sprints 26.11–26.13", bug_bullets,
                  MARGIN, PptInches(4.05), PptInches(12.1), PptInches(2.5), CORAL, font_size=8.5)
    slide_footer(slide, ui["maintenance_note"], n := n + 1)

    # 8 Perf — 3-sprint deliverables only
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Performance Testing", perf["scope_note"], "Perf")
    deliverables = perf["sprint_deliverables"]
    for i, d in enumerate(deliverables):
        add_kpi_card(slide, MARGIN + PptInches(i * 4.0), PptInches(1.55), PptInches(3.7), PptInches(1.35),
                     f"Sprint {d['sprint']}", d["name"][:40] + ("…" if len(d["name"]) > 40 else ""), CORAL, d["jira"])
    slide_bullets(slide, "Delivered (last 3 sprints)", [f"{d['name']}: {d['status']}" for d in deliverables],
                  MARGIN, PptInches(3.15), PptInches(12.1), PptInches(2.0), TEAL, font_size=9)
    slide_bullets(slide, "Ongoing baseline track", [perf["ongoing"]],
                  MARGIN, PptInches(5.35), PptInches(12.1), PptInches(1.2), NAVY, font_size=9)
    slide_footer(slide, "Emergency work delivered — governance ask on next slide", n := n + 1)

    # 9 Perf governance ask
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Performance Governance — Leadership Ask", "Stop last-minute emergency intake", "Ask")
    slide_bullets(slide, "What happened", [
        "Two emergency perf requests in two sprints — barcode (Synergy) and Jahia/IDP server patch",
        "Teams approached us days before deploy — we dropped sprint work and delivered in ~1 week each",
        "Scripts ran and results accepted — but KB, check-in, test data, and standards steps were skipped",
    ], MARGIN, PptInches(1.5), PptInches(5.9), PptInches(2.6), CORAL)
    slide_bullets(slide, "What we need from leadership", [
        "Add QA Automation at SDLC design — not at deploy approval",
        "Mandatory perf gate for IDP login, enrollment, contributions, withdrawals before release",
        "Route requests through standard intake — time to document, check in, and expand coverage",
        "We will always help in emergencies — but repeated ad-hoc work blocks the roadmap",
    ], PptInches(6.7), PptInches(1.5), PptInches(5.9), PptInches(2.6), TEAL)
    slide_footer(slide, "Delivered under pressure — asking for a sustainable model", n := n + 1)

    # 10 Capacity
    slide_chart_page(prs, n := n + 1, "Capacity — MSC Wrap-up vs New Work",
        "1 resource close-out · 3.5 FTE for next priority", "03-capacity-timeline.png",
        [
            cap["msc_wrap"],
            f"New work capacity: {cap['new_work']}",
            "Sunil → MSC wrap-up: qTest, Bruno, coverage matrix, handover docs",
            "Preeti → perf: contribution JMX now; M1+enrollment perf next sprint",
            "Main squad → ACS-5678 Atlas (Oct) OR V2/V3 backlog — need leadership queue",
        ], "Coding complete end 26.14 — not a full squad exit from MSC", "Capacity")

    # 11 V2 sunset
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "V2 → V3 Migration Path", "Reduce Jenkins footprint", "Strategy")
    vs = m["v2_sunset"]
    slide_bullets(slide, "Plan", [
        vs["goal"],
        vs["driver"],
        vs["discussion"],
        f"Lead: {vs['lead']}",
        "Outcome: lighter GitLab-based V3 suites; retire duplicate V2 enrollment from Jenkins",
    ], MARGIN, PptInches(1.5), PptInches(12.1), PptInches(4.5), VIOLET, font_size=9.5)
    slide_footer(slide, "Aligns with platform direction — V3 as source of truth", n := n + 1)

    # 12 ACS-5678
    inc = m["incoming_work"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Incoming: ACS-5678 Atlas", "Traffic distribution & load balancer rebalancing", "New work")
    slide_bullets(slide, "AHA summary (Brian Danilczyk)", [
        f"{inc['aha']} — {inc['priority']}",
        f"Estimate: {inc['estimate']} · Target: {inc['target']}",
        f"Status: {inc.get('status', 'Intake pending')}",
        f"Dependency: {inc['dependency']}",
        "Scope: validate traffic distribution after Redis session externalization; load + perf testing",
        "Main squad can start intake next sprint while 1 resource finishes MSC wrap-up",
    ], MARGIN, PptInches(1.5), PptInches(12.1), PptInches(4.2), CYAN)
    slide_footer(slide, "Confirm priority vs MSC wrap-up, V2 sunset, and other AHA items", n := n + 1)

    # 13 Leadership asks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    slide_header(slide, "Leadership Asks & Decisions Needed", "Please respond on this call — by area", "Decisions")
    slide_categorized_asks(slide, m["leadership_asks"], MARGIN, PptInches(1.45), PptInches(12.1), PptInches(5.35))
    slide_footer(slide, "Full ask list in companion DOCX · Swapnil Patil · QA Automation AM Squad", n := n + 1)

    # 14 Close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, dark=True)
    t = slide.shapes.add_textbox(MARGIN, PptInches(3.0), PptInches(11), PptInches(1))
    t.text_frame.paragraphs[0].text = "Questions & Discussion"
    t.text_frame.paragraphs[0].font.size = PptPt(36)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    s = slide.shapes.add_textbox(MARGIN, PptInches(4.1), PptInches(11), PptInches(0.5))
    s.text_frame.paragraphs[0].text = "Detailed briefing document shared · KB: programs/leadership-updates/2026-08-28-rajib-henry-biweekly/"
    s.text_frame.paragraphs[0].font.size = PptPt(12)
    s.text_frame.paragraphs[0].font.color.rgb = TEAL

    out = DELIVERABLES / "AM-Squad-Biweekly-Status-Rajib-Henry-Aug28-2026.pptx"
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


# ── DOCX helpers ─────────────────────────────────────────────────────────────

def shade_cell(cell, fill_hex: str) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill_hex)
    tc_pr.append(shd)


def styled(doc, text, bold=False, size=11, color=None, italic=False):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = bold
    r.italic = italic
    r.font.size = Pt(size)
    if color:
        r.font.color.rgb = color
    return p


def add_table(doc, headers, rows, header_fill="003241"):
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        c.text = h
        shade_cell(c, header_fill)
        for r in c.paragraphs[0].runs:
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            r.font.size = Pt(9)
    for row in rows:
        cells = t.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = str(val)
            for r in cells[i].paragraphs[0].runs:
                r.font.size = Pt(9)
    doc.add_paragraph()


def build_docx(m: dict) -> Path:
    doc = Document()
    doc.styles["Normal"].font.name = "Segoe UI"
    doc.styles["Normal"].font.size = Pt(11)
    meeting = datetime.strptime(m["meeting_date"], "%Y-%m-%d").strftime("%B %d, %Y")

    # Cover
    styled(doc, "QA Automation — AM Squad", bold=True, size=24, color=DOCX_NAVY)
    styled(doc, f"Biweekly Status Briefing · {meeting}", size=14, color=DOCX_TEAL)
    styled(doc, "Audience: Rajib (Chapter Lead) · Henry (Director) · Persistent Delivery Managers", size=11, color=DOCX_GRAY)
    doc.add_paragraph()
    meta = doc.add_table(rows=4, cols=2)
    meta.style = "Table Grid"
    rows_meta = [
        ("Prepared by", "Swapnil Patil / QA Automation AM Squad"),
        ("Sprints covered", m["sprints"]["scope"]),
        ("Classification", "Internal — Leadership"),
        ("Companion deck", "AM-Squad-Biweekly-Status-Rajib-Henry-Aug28-2026.pptx"),
    ]
    for i, (label, val) in enumerate(rows_meta):
        meta.rows[i].cells[0].text = label
        meta.rows[i].cells[1].text = val
    doc.add_page_break()

    # Executive summary
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(
        "The AM Squad continues parallel delivery across Unite MSC API automation, V2/V3 UI regression, "
        "performance testing, and pipeline/standards work. This briefing covers Sprints 26.13 and "
        "26.14 (current). Mobile 1 and Mobile 2 are 100% of in-scope business endpoints (26 + 25) and "
        "ready for sign-off. Enrollment core E2E is ~95% (19/20 endpoints coded; 17 TestNG classes). "
        "Submit (review-confirm-entered) is the remaining coding step — target end Sprint 26.14. "
        "Coding complete is not program closed: one resource stays ~1 sprint on documentation, qTest, "
        "Bruno handover, and sign-off packaging while the main squad pivots to Atlas or V2/V3 backlog."
    )
    p = m["pulse"]
    u = m["unite_msc"]
    add_table(doc, ["Metric", "Value", "Notes"], [
        ("Sprint scope", p["sprint_scope"], "26.13 · 26.14"),
        ("GitLab merges", p["gitlab_merges"], p["gitlab_merges_note"]),
        ("Sprint 26.14 board items", p["jira_sprint_items_current"], p["jira_sprint_items_note"]),
        ("Regression bugs", p["automation_bugs_logged"], p["automation_bugs_note"]),
        ("MSC in-scope APIs", p["msc_in_scope_endpoints"], p["msc_in_scope_note"]),
        ("MSC Mobile 2", f"{u['mobile2']['automated']}/{u['mobile2']['total']}", u["mobile2"]["status"]),
        ("MSC Mobile 1", f"{u['mobile1']['automated']}/{u['mobile1']['total']}", u["mobile1"]["status"]),
        ("MSC Enrollment", f"~{u['enrollment']['pct']}% ({u['enrollment']['core_e2e']})", u["enrollment"]["status"]),
        ("Delivery vs ETA", p["eta_delivery_rate"], "Consistent ahead-of-plan delivery"),
    ])
    doc.add_page_break()

    # MSC detail
    doc.add_heading("1. Unite MSC API Automation", level=1)
    doc.add_paragraph(
        "Program KB: programs/unite-msc/ · Repository: api-test-automation (GitLab). "
        "Delivered in approximately 50% of original ETA using AI-assisted migration and canonical TestNG framework."
    )
    doc.add_heading("Coverage snapshot", level=2)
    add_table(doc, ["Module", "Coverage", "Status"], [
        ("Mobile 2", "25/25 endpoints (100%)", "Complete — sign-off package ready (QA-1553, QA-1567)"),
        ("Mobile 1", "26/26 in-scope operations (100%)", u["mobile1"]["status"]),
        ("Enrollment", f"{u['enrollment']['core_e2e']} core E2E ({u['enrollment']['pct']}%)", u["enrollment"]["status"]),
    ])
    doc.add_paragraph(u["footnote"])
    doc.add_heading("Pipeline & environments", level=2)
    for item in [
        u["pipeline"]["stage1"],
        u["pipeline"]["qc4"],
        u["pipeline"]["gitlab_nightly"],
        "Hybrid model per Jul 2025 leadership decision — Stage1 is primary until QC4 stabilizes.",
        "Blocker: no non-IDP plants in QC4 lower environment — need product owner to name test plans.",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Sprint 26.13–26.14 MSC highlights (Jira)", level=2)
    highlights = [
        "QA-1786–1789 — Enrollment optional helpers: mobile login, enrollment started, recurring contribution, allocation funds",
        "QA-1775–1776, QA-1602–1603 — Wizard steps: owner address, content, allocations, verify routing",
        "QA-1405 — Mobile 2 GitLab nightly regression job delivered",
        "QA-1668–1671 — Mobile 1/2 coverage validation and sign-off prep complete",
        "QA-1604 / QA-1810 — Submit (review-confirm-entered) — in flight end Sprint 26.14",
        "DB refresh: test data + SQL stabilized for Stage1/QC4 enrollment runs (Aug 2026)",
    ]
    for h in highlights:
        doc.add_paragraph(h, style="List Bullet")
    if (ASSETS / "04-msc-endpoint-coverage.png").exists():
        doc.add_paragraph()
        doc.add_picture(str(ASSETS / "04-msc-endpoint-coverage.png"), width=Inches(5.5))
    doc.add_page_break()

    # V2/V3
    doc.add_heading("2. V2 & V3 UI Automation", level=1)
    ui = m["ui_automation"]
    doc.add_paragraph(
        f"{ui['scope_note']}. {ui['maintenance_note']}."
    )
    doc.add_heading("Sprint 26.11–26.13 highlights", level=2)
    for item in ui["v3_highlights"] + ui["v2_highlights"]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("V2 → V3 migration — discussion needed", level=2)
    doc.add_paragraph(ui["migration_discussion"])
    doc.add_heading("Regression bugs (Sprints 26.11–26.13)", level=2)
    add_table(doc, ["Key", "Summary", "Area", "Sprint"], [[b["key"], b["summary"], b["area"], b["sprint"]] for b in m["recent_bugs"]])
    doc.add_page_break()

    # Performance
    doc.add_heading("3. Performance Testing", level=1)
    perf = m["performance"]
    doc.add_paragraph(perf["scope_note"] + ". " + perf["ongoing"])
    doc.add_heading("Deliverables (Sprints 26.11–26.13)", level=2)
    add_table(doc, ["Request", "Sprint", "Status", "Jira"], [
        (e["name"], e["sprint"], e["status"], e.get("jira", "")) for e in perf["sprint_deliverables"]
    ])
    doc.add_heading("Performance governance ask", level=2)
    doc.add_paragraph(
        "Two emergency requests in two sprints forced the team to drop planned sprint work. We delivered "
        "scripts and baseline results in approximately one week each, but the full standard — KB documentation, "
        "Git check-in, test data management, and expanded regression coverage — was not possible under "
        "last-minute intake."
    )
    for item in [
        "Engage QA Automation when performance requirements are designed — not at deploy approval",
        "Mandate performance validation for IDP login, enrollment, contributions, and withdrawals before release",
        "Route cross-team perf requests through standard intake with lead time",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # Capacity & roadmap
    doc.add_heading("4. Capacity & Upcoming Work", level=1)
    cap = m["capacity"]
    doc.add_paragraph(
        f"MSC program target: {u['target_completion']}. "
        f"Expected: {cap['msc_wrap']}."
    )
    doc.add_heading("MSC wrap-up checklist (1 resource)", level=2)
    for item in cap.get("closeout_checklist", []):
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Planned reallocation", level=2)
    for item in cap["reallocation"]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("V2 → V3 migration", level=2)
    vs = m["v2_sunset"]
    doc.add_paragraph(vs["goal"])
    doc.add_paragraph(vs["driver"])
    doc.add_paragraph(f"Lead: {vs['lead']}")
    doc.add_heading("Incoming: ACS-5678 Atlas (AHA)", level=2)
    inc = m["incoming_work"]
    for item in [
        f"Feature: {inc['aha']}",
        f"Owner: {inc['owner_aha']} · Priority: {inc['priority']}",
        f"Estimate: {inc['estimate']} · Dependency: {inc['dependency']}",
        f"Target: {inc['target']}",
        "Scope: load/perf/regression validation of traffic distribution after Redis session externalization",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # Leadership asks
    doc.add_heading("5. Leadership Asks & Decisions", level=1)
    for category, asks in m["leadership_asks"].items():
        doc.add_heading(category, level=2)
        for ask in asks:
            doc.add_paragraph(ask, style="List Bullet")
    doc.add_heading("Sign-off owners needed (names TBD from leadership)", level=2)
    add_table(doc, ["Area", "What we need", "Suggested role"], [
        ("Mobile 2 API", "Green signal for 25/25 endpoint coverage", "MSC API SME / product owner"),
        ("Mobile 1 API", "Accept 26/26 in-scope endpoints (suite-tier routing)", "MSC API SME"),
        ("MSC Pipeline", "Approve GitLab nightly + ENVP integration", "DevOps + program lead"),
        ("Performance baselines", "Accept IDP/MSC perf thresholds", "Platform / perf SME"),
        ("Non-IDP test plans", "Name plans for QC4 lower env", "Product / plan configuration owner"),
    ])
    doc.add_page_break()

    doc.add_heading("Appendix — References", level=1)
    refs = [
        "programs/unite-msc/leadership/2026-08-07-kevin-status-update.md",
        "programs/unite-msc/leadership/2026-07-17-leadership-update/",
        "programs/unite-msc/leadership/2026-07-23-scope-alignment/",
        "programs/leadership-updates/2026-08-am-squad-leadership-update/",
        "programs/Performance Testing/barcode-syn-443/",
        "Jira: QA project · AMSQUAD sprints",
    ]
    for r in refs:
        doc.add_paragraph(r, style="List Bullet")
    styled(doc, f"Generated {datetime.now().strftime('%B %d, %Y')}.", italic=True, size=9, color=DOCX_GRAY)

    out = DELIVERABLES / "AM-Squad-Biweekly-Status-Rajib-Henry-Aug28-2026.docx"
    try:
        doc.save(out)
    except PermissionError:
        out = DELIVERABLES / "AM-Squad-Biweekly-Status-Rajib-Henry-Aug28-2026-UPDATED.docx"
        doc.save(out)
        print(f"Note: original DOCX locked — wrote {out.name} instead")
    return out


def main() -> None:
    m = load_data()
    generate_charts(m)
    ppt = build_pptx(m)
    doc = build_docx(m)
    print(f"Wrote {ppt}")
    print(f"Wrote {doc}")


if __name__ == "__main__":
    main()
