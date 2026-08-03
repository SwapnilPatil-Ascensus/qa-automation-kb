#!/usr/bin/env python3
"""Generate Automation Bug Lifecycle standard DOCX + PPTX deliverables."""

from __future__ import annotations

from pathlib import Path

import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from PIL import Image

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
DOCX_OUT = ROOT / "deliverables" / "Automation-Bug-Lifecycle-Playbook.docx"
PPTX_OUT = ROOT / "deliverables" / "Automation-Bug-Lifecycle-Standard.pptx"
VERSION_DATE = "July 24, 2026"

# Ascensus-aligned palette
NOBLE = "#003241"
TEAL = "#026B84"
PEAK = "#009E86"
SKY = "#05A2C6"
ACTION = "#2D65B4"
SUCCESS = "#1AA01C"
WARNING = "#C27800"
ALERT = "#D2260F"
SLATE = "#475569"
MUTED = "#94A3B8"
LIGHT = "#F1F5F9"
WHITE = "#FFFFFF"
RED_BG = "#FFEBE8"
AMBER_BG = "#FFF4E5"
GREEN_BG = "#E8F8E8"
LIGHT_BG_CHART = "#E8F4F8"

NAVY = RGBColor(0x00, 0x32, 0x41)
TEAL_RGB = RGBColor(0x02, 0x6B, 0x84)
GREEN_RGB = RGBColor(0x1A, 0xA0, 0x1C)
AMBER_RGB = RGBColor(0xC2, 0x78, 0x00)
RED_RGB = RGBColor(0xD2, 0x26, 0x0F)
GRAY_RGB = RGBColor(0x47, 0x55, 0x69)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
DOC_LIGHT_BG = "E8F4F8"
DOC_GREEN_BG = "E8F8E8"
DOC_AMBER_BG = "FFF4E5"
DOC_RED_BG = "FFEBE8"

PPT_NOBLE = PptRGB(0x00, 0x32, 0x41)
PPT_NOBLE_DARK = PptRGB(0x00, 0x1A, 0x22)
PPT_TEAL = PptRGB(0x05, 0xA2, 0xC6)
PPT_TEAL_DARK = PptRGB(0x02, 0x6B, 0x84)
PPT_PEAK = PptRGB(0x00, 0x9E, 0x86)
PPT_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
PPT_OFF = PptRGB(0xF8, 0xFA, 0xFC)
PPT_SLATE = PptRGB(0x64, 0x74, 0x8B)
PPT_TEXT = PptRGB(0x1E, 0x29, 0x3B)
PPT_MUTED = PptRGB(0x94, 0xA3, 0xB8)
PPT_SUCCESS = PptRGB(0x1A, 0xA0, 0x1C)
PPT_SUCCESS_BG = PptRGB(0xE8, 0xF8, 0xE8)
PPT_WARN = PptRGB(0xC2, 0x78, 0x00)
PPT_WARN_BG = PptRGB(0xFF, 0xF4, 0xE5)
PPT_ALERT = PptRGB(0xD2, 0x26, 0x0F)
PPT_ALERT_BG = PptRGB(0xFF, 0xEB, 0xE8)
PPT_ACTION = PptRGB(0x2D, 0x65, 0xB4)
PPT_CARD = PptRGB(0xFF, 0xFF, 0xFF)
PPT_BORDER = PptRGB(0xE2, 0xE8, 0xF0)

SLIDE_W = PptInches(13.333)
SLIDE_H = PptInches(7.5)
MARGIN = PptInches(0.6)
HEADER_H = PptInches(1.1)
BODY_TOP = PptInches(1.42)
FOOTER_Y = PptInches(7.08)

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Segoe UI", "Calibri", "Arial"],
    "figure.facecolor": WHITE,
})


def _rounded_box(ax, x, y, w, h, face, edge, alpha=1.0, lw=2, radius=0.08):
    box = FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0.02,rounding_size={radius}",
        facecolor=face, edgecolor=edge, linewidth=lw, alpha=alpha, zorder=2,
    )
    ax.add_patch(box)
    return box


def _arrow(ax, x1, y1, x2, y2, color=SLATE, lw=2.2, rad=0.0):
    style = f"arc3,rad={rad}" if rad else "arc3,rad=0.0"
    arr = FancyArrowPatch(
        (x1, y1), (x2, y2),
        arrowstyle="-|>", mutation_scale=14, linewidth=lw,
        color=color, connectionstyle=style, zorder=1,
        shrinkA=4, shrinkB=4,
    )
    ax.add_patch(arr)


def _center_text(ax, x, y, text, size=11, color=NOBLE, bold=False, lines=None):
    if lines:
        for i, line in enumerate(lines):
            ax.text(x, y - i * 0.28, line, ha="center", va="center",
                    fontsize=size, fontweight="bold" if bold else "normal", color=color)
    else:
        ax.text(x, y, text, ha="center", va="center",
                fontsize=size, fontweight="bold" if bold else "normal", color=color)


def chart_triage_decision(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12, 7))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 8)
    ax.axis("off")
    fig.patch.set_facecolor(WHITE)

    # Root
    _rounded_box(ax, 4.6, 6.55, 2.8, 0.95, NOBLE, NOBLE)
    _center_text(ax, 6.0, 7.02, "Failure Detected", size=13, color=WHITE, bold=True)

    branches = [
        (1.2, 3.6, 2.6, 1.35, RED_BG, ALERT,
         ["Environment", "DB refresh · certs · OKD"],
         ["No JIRA", "No main lock"], 0.12),
        (4.7, 3.6, 2.6, 1.35, AMBER_BG, WARNING,
         ["Flaky / False failure", "Pass on retry · timing"],
         ["Flakiness playbook", "Quarantine if needed"], 0.0),
        (8.2, 3.6, 2.6, 1.35, GREEN_BG, SUCCESS,
         ["Functional defect", "Reproducible · app trace"],
         ["Full bug cycle", "JIRA · notify · lock"], -0.12),
    ]

    for x, y, w, h, face, edge, title_lines, action_lines, rad in branches:
        _rounded_box(ax, x, y, w, h, face, edge, lw=2.2)
        cy = y + h / 2 + 0.12
        for i, line in enumerate(title_lines):
            _center_text(ax, x + w / 2, cy - i * 0.32, line, size=11, color=edge, bold=(i == 0))
        ay = y - 0.55
        for i, line in enumerate(action_lines):
            _center_text(ax, x + w / 2, ay - i * 0.32, line, size=9.5, color=SLATE)
        # Arrow: bottom of root → top center of branch
        _arrow(ax, 6.0, 6.55, x + w / 2, y + h, color=SLATE, rad=rad)

    ax.text(6.0, 0.35,
            "Validate locally before logging. Critical legitimate failures may trigger monolith/automation main lock (10 AM SLA).",
            ha="center", fontsize=9, color=MUTED, style="italic")
    ax.text(6.0, 7.75, "Triage Decision Tree — Bug or Not?",
            ha="center", fontsize=17, fontweight="bold", color=NOBLE)
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.35)
    plt.close(fig)


def chart_end_to_end_workflow(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14, 5.5))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 5.5)
    ax.axis("off")

    steps = [
        (0.15, "Detect", ALERT),
        (1.95, "Triage", WARNING),
        (3.75, "Evidence", TEAL),
        (5.55, "Cursor H", ACTION),
        (7.35, "JIRA +\nNotify", SUCCESS),
        (9.15, "Change\nSet", NOBLE),
        (10.95, "Resolve", PEAK),
    ]
    bw, bh = 1.55, 1.05
    y = 3.0
    for i, (x, label, color) in enumerate(steps):
        _rounded_box(ax, x, y, bw, bh, color, color, alpha=0.14, lw=2.5)
        _center_text(ax, x + bw / 2, y + bh / 2, label, size=10.5, color=color, bold=True)
        if i < len(steps) - 1:
            nx = steps[i + 1][0]
            _arrow(ax, x + bw, y + bh / 2, nx, y + bh / 2, color=SLATE, lw=2.5)

    notes = [
        (1.95, 1.55, "Env / flaky → fix locally", SUCCESS, GREEN_BG),
        (5.55, 1.55, "Leadership-approved templates", ACTION, LIGHT_BG_CHART),
        (9.15, 1.55, "GitLab PM: MRs + commits", TEAL, LIGHT),
    ]
    for x, y_n, text, color, bg in notes:
        _rounded_box(ax, x, y_n, 2.8, 0.72, bg, color, lw=1.5, alpha=0.95)
        _center_text(ax, x + 1.4, y_n + 0.36, text, size=9, color=color)

    ax.text(7.0, 5.05, "Automation Regression Failure — End-to-End Standard Workflow",
            ha="center", fontsize=16, fontweight="bold", color=NOBLE)
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.3)
    plt.close(fig)


def chart_toolchain(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12, 4.8))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4.8)
    ax.axis("off")

    tools = [
        (0.4, "qa-automation-kb\nCursor · Prompt H", NOBLE, "JIRA · Email · Teams"),
        (3.2, "JIRA\nQA Board", TEAL, "Ticket · RCA · Status"),
        (6.0, "GitLab Project\nManager", PEAK, "MRs · Commits · Authors"),
        (8.8, "Stakeholders\nTeams · Email", ACTION, "Leadership-approved"),
    ]
    for x, title, color, out in tools:
        _rounded_box(ax, x, 2.4, 2.4, 1.35, color, color, alpha=0.12, lw=2.2)
        for i, line in enumerate(title.split("\n")):
            _center_text(ax, x + 1.2, 3.15 - i * 0.35, line, size=10.5, color=color, bold=(i == 0))
        _rounded_box(ax, x + 0.15, 1.35, 2.1, 0.65, LIGHT, MUTED, lw=1)
        _center_text(ax, x + 1.2, 1.67, out, size=8.5, color=SLATE)

    for x1, x2 in [(2.8, 3.2), (5.6, 6.0), (8.4, 8.8)]:
        _arrow(ax, x1, 3.08, x2, 3.08, color=SLATE)

    ax.text(6.0, 4.35, "Integrated Toolchain — Evidence to Footprint",
            ha="center", fontsize=16, fontweight="bold", color=NOBLE)
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.3)
    plt.close(fig)


def chart_defect_lifecycle(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12, 2.8))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 2.8)
    ax.axis("off")
    states = ["New", "Assigned", "In Progress", "Fixed", "Verified", "Closed"]
    colors = [NOBLE, TEAL, TEAL, SUCCESS, SUCCESS, SLATE]
    sw = 1.65
    gap = 0.22
    y = 1.0
    for i, (state, color) in enumerate(zip(states, colors)):
        x = 0.35 + i * (sw + gap)
        _rounded_box(ax, x, y, sw, 0.85, color, color, alpha=0.18, lw=2)
        _center_text(ax, x + sw / 2, y + 0.42, state, size=11, color=color, bold=True)
        if i < len(states) - 1:
            _arrow(ax, x + sw, y + 0.42, x + sw + gap, y + 0.42, color=SLATE, lw=2)
    ax.text(6.0, 2.35, "Defect Lifecycle — JIRA Workflow", ha="center", fontsize=15, fontweight="bold", color=NOBLE)
    ax.text(6.0, 0.35, "Branches: Duplicate · Rejected · Deferred · Reopened",
            ha="center", fontsize=9, color=MUTED, style="italic")
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.25)
    plt.close(fig)


def chart_repo_priority(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 4.5))
    repos = ["monolith (application)", "automation", "qa-automation", "prime-test-automation"]
    scores = [95, 85, 72, 68]
    colors = [NOBLE, TEAL, PEAK, SLATE]
    y_pos = range(len(repos))
    bars = ax.barh(y_pos, scores, color=colors, height=0.55, edgecolor=WHITE, linewidth=2)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(repos, fontsize=11)
    ax.set_xlim(0, 108)
    ax.set_xlabel("Investigation priority when application-area failure is suspected", fontsize=10, color=SLATE)
    ax.set_title("GitLab Project Manager — Repository Investigation Order",
                 fontsize=14, fontweight="bold", color=NOBLE, pad=14)
    for bar, score in zip(bars, scores):
        ax.text(score + 1.5, bar.get_y() + bar.get_height() / 2, f"{score}%",
                va="center", fontsize=11, fontweight="bold", color=NOBLE)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CBD5E1")
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.tick_params(colors=SLATE)
    fig.text(0.5, 0.01, "Query merged MRs and commits between last green run and failure. Export or copy into JIRA / email.",
             ha="center", fontsize=8.5, color=MUTED)
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.3)
    plt.close(fig)


def chart_evidence_folder(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")
    _rounded_box(ax, 0.5, 0.5, 9.0, 4.2, LIGHT, TEAL, lw=2)
    ax.text(5.0, 4.35, "automation-bug-lifecycle/evidence/regression-reports/[MMDDYYYY]/", ha="center",
            fontsize=12, fontweight="bold", color=NOBLE, family="monospace")
    items = [
        ("Screenshots", ".png — failure UI captures"),
        ("Exception logs", ".txt — stack traces from CI"),
        ("Test data", "Test Data.txt · testcase_information"),
        ("Console output", "GitLab / Jenkins job log"),
        ("Bug documentation", "[MMDDYYYY]_[Feature]_[Issue].md"),
        ("TestNG report", "URL or exported HTML index"),
    ]
    for i, (title, desc) in enumerate(items):
        y = 3.55 - i * 0.55
        _rounded_box(ax, 0.85, y - 0.18, 8.3, 0.48, WHITE, "#CBD5E1", lw=1.2)
        ax.text(1.1, y + 0.05, title, fontsize=10.5, fontweight="bold", color=TEAL)
        ax.text(3.2, y + 0.05, desc, fontsize=9.5, color=SLATE)
    ax.text(5.0, 4.75, "Evidence Folder — Standard Structure", ha="center",
            fontsize=15, fontweight="bold", color=NOBLE)
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor=WHITE, pad_inches=0.3)
    plt.close(fig)


def generate_charts() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    charts = {
        "triage": ASSETS / "chart_triage_decision.png",
        "workflow": ASSETS / "chart_end_to_end_workflow.png",
        "toolchain": ASSETS / "chart_toolchain.png",
        "lifecycle": ASSETS / "chart_defect_lifecycle.png",
        "repos": ASSETS / "chart_repo_priority.png",
        "evidence": ASSETS / "chart_evidence_folder.png",
    }
    chart_triage_decision(charts["triage"])
    chart_end_to_end_workflow(charts["workflow"])
    chart_toolchain(charts["toolchain"])
    chart_defect_lifecycle(charts["lifecycle"])
    chart_repo_priority(charts["repos"])
    chart_evidence_folder(charts["evidence"])
    return charts


# ── DOCX helpers ──────────────────────────────────────────────────────────────

def shade_cell(cell, hex_color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False, color: RGBColor | None = None, size: int = 10) -> None:
    cell.text = ""
    run = cell.paragraphs[0].add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color


def style_table_header(row) -> None:
    for cell in row.cells:
        shade_cell(cell, "003241")
        for p in cell.paragraphs:
            for run in p.runs:
                run.font.color.rgb = WHITE_RGB
                run.bold = True
                run.font.size = Pt(10)
                run.font.name = "Calibri"


def add_heading(doc: Document, title: str, level: int = 1) -> None:
    h = doc.add_heading(title, level=level)
    for run in h.runs:
        run.font.color.rgb = NAVY
        run.font.name = "Calibri"


def add_callout(doc: Document, text: str, bg: str = DOC_LIGHT_BG) -> None:
    table = doc.add_table(rows=1, cols=1)
    cell = table.rows[0].cells[0]
    shade_cell(cell, bg)
    set_cell_text(cell, text, size=11)
    for p in cell.paragraphs:
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        for run in p.runs:
            run.font.color.rgb = NAVY


def setup_header_footer(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    hp = section.header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.add_run("QA Automation Standard  |  Automation Bug Lifecycle  |  Internal")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY_RGB
    run.font.name = "Calibri"
    fp = section.footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run(f"QA Automation — AMSQUAD  |  Version {VERSION_DATE}")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY_RGB
    run.font.name = "Calibri"


def build_docx(charts: dict[str, Path]) -> None:
    DOCX_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    setup_header_footer(doc)

    for _ in range(5):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("AUTOMATION BUG LIFECYCLE")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = NAVY
    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = s.add_run("QA Automation Operating Standard")
    r2.font.size = Pt(20)
    r2.font.color.rgb = TEAL_RGB
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for line in [f"Version {VERSION_DATE}", "Applicable to all QA Automation teams and programs",
                 "Single source of truth: qa-automation-kb repository"]:
        run = meta.add_run(line + "\n")
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_RGB
    doc.add_page_break()

    add_heading(doc, "1. Purpose & Scope", level=1)
    add_callout(doc,
        "This standard defines how QA Automation responds when regression or release tests fail. "
        "It applies to any team running V2 (Jenkins) or V3 (GitLab) suites, API automation, "
        "or performance regression — not tied to a single program or client.")
    doc.add_paragraph(
        "The workflow integrates: (1) triage per automation-bug-lifecycle/process/TRIAGE_RULES.md and FLAKINESS_PLAYBOOK.md, "
        "(2) evidence collection in automation-bug-lifecycle/evidence/regression-reports/, (3) Cursor Prompt H from qa-knowledge-base/00_SYSTEM/PROMPTS.md, "
        "(4) leadership-approved communication templates from Confluence Bug Handling exports, and "
        "(5) GitLab Project Manager for change-set investigation.")
    doc.add_page_break()

    add_heading(doc, "2. Triage — Bug or Not?", level=1)
    doc.add_picture(str(charts["triage"]), width=Inches(6.8))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    triage = doc.add_table(rows=5, cols=4)
    triage.style = "Table Grid"
    for i, h in enumerate(["Classification", "Signals", "Action", "JIRA?"]):
        set_cell_text(triage.rows[0].cells[i], h, bold=True, color=WHITE_RGB)
    style_table_header(triage.rows[0])
    for ri, row in enumerate([
        ("Environment", "DB restricted, certs/helm, OKD down, env-only", "Escalate infra", "No"),
        ("Flaky / false failure", "Pass on retry; timing; locator; TR flip same day", "Flakiness playbook", "No"),
        ("Automation script", "Test logic, wait, data setup", "Fix locally in automation repo", "Optional"),
        ("Functional defect", "Reproducible; manual repro; app stack trace", "Full bug cycle", "Yes"),
    ], start=1):
        for ci, val in enumerate(row):
            set_cell_text(triage.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY_RGB)
    doc.add_page_break()

    add_heading(doc, "3. End-to-End Workflow", level=1)
    doc.add_picture(str(charts["workflow"]), width=Inches(6.8))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    steps = [
        ("Detect", "V2 Jenkins or V3 GitLab regression reports failure. Record last green run timestamp."),
        ("Triage", "Classify per Section 2. Rerun locally. Check environment health (automation-bug-lifecycle/process/DAILY_REGRESSION.md)."),
        ("Evidence", "Folder: automation-bug-lifecycle/evidence/regression-reports/[MMDDYYYY]/ — screenshots, logs, test data."),
        ("Cursor H", "qa-knowledge-base/00_SYSTEM/PROMPTS.md section H → JIRA block + email + Teams in one .md file."),
        ("JIRA + Notify", "QA board ticket. Leadership-approved To/Cc. Critical legit → lock main (10 AM SLA)."),
        ("Change set", "GitLab Project Manager: MRs + commits between last pass and failure."),
        ("Resolve", "Fix verified → JIRA closed with RCA → resolution email → unlock main."),
    ]
    st = doc.add_table(rows=len(steps) + 1, cols=3)
    st.style = "Table Grid"
    for i, h in enumerate(["Step", "Activity", "Standard reference"]):
        set_cell_text(st.rows[0].cells[i], h, bold=True, color=WHITE_RGB)
    style_table_header(st.rows[0])
    refs = ["automation-bug-lifecycle/process/DAILY_REGRESSION.md", "FLAKINESS_PLAYBOOK.md", "BUG_REPORTING_PROCESS.md",
            "qa-knowledge-base/00_SYSTEM/PROMPTS.md", "Confluence Bug Handling PDFs", "GitLab Project Manager", "1b Resolution PDF"]
    for ri, ((act, det), ref) in enumerate(zip(steps, refs), start=1):
        set_cell_text(st.rows[ri].cells[0], str(ri), bold=True, color=NAVY)
        set_cell_text(st.rows[ri].cells[1], f"{act}: {det}")
        set_cell_text(st.rows[ri].cells[2], ref, size=9)
    doc.add_page_break()

    add_heading(doc, "4. Evidence Folder Structure", level=1)
    doc.add_picture(str(charts["evidence"]), width=Inches(6.2))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("Naming: [MMDDYYYY]_[FeatureName]_[IssueType].md — see BUG_REPORTING_PROCESS.md.")

    add_heading(doc, "5. Cursor Prompt H Deliverables", level=1)
    doc.add_picture(str(charts["toolchain"]), width=Inches(6.2))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_callout(doc,
        "Teams and email templates produced by Prompt H follow formats approved by leadership and senior QA resources. "
        "Use standard To/Cc lists — do not improvise distribution or structure.", DOC_GREEN_BG)

    add_heading(doc, "6. GitLab Project Manager — Change Set", level=1)
    mr_img = ASSETS / "gitlab_mr_results_sample.png"
    ui_img = ASSETS / "gitlab_project_manager_ui.png"
    if ui_img.exists():
        doc.add_picture(str(ui_img), width=Inches(5.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    if mr_img.exists():
        doc.add_picture(str(mr_img), width=Inches(5.8))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap = doc.add_paragraph("Sample output: Merge Requests for monolith — date range, authors, merged by, branches")
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_picture(str(charts["repos"]), width=Inches(5.5))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    for item in [
        "Utility: C:\\Development\\Workspace\\GitlabInfoProjUI (separate repo; read-only GitLab API)",
        "Tabs: Merge Requests · Commits · Branches · Pipelines · Users · Project",
        "Filter Merged MRs; note Author, Merged by, source → target branch, Merged At (EST)",
        "Export CSV/Excel or copy formatted report into JIRA / email Change Set section",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    add_heading(doc, "7. Defect Lifecycle & Branch Locking", level=1)
    doc.add_picture(str(charts["lifecycle"]), width=Inches(6.5))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("Full reference: automation-bug-lifecycle/process/DEFECT_LIFECYCLE.md, automation-bug-lifecycle/process/TRIAGE_RULES.md")
    add_callout(doc, "Lock monolith/main + automation/main only for critical, legitimate defects — not env/data/flaky.", DOC_AMBER_BG)

    add_heading(doc, "8. Multi-Failure Rollup", level=1)
    for item in [
        "Group by feature/plan — not one JIRA per test method",
        "Failure matrix with root-cause hints per traunch/plan",
        "Umbrella ticket + linked children when appropriate",
        "Reference: automation-bug-lifecycle/evidence/regression-reports/04202026/04202026_DailyRegression_PipelineFailureRollup.md",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "9. Repository Cross-References", level=1)
    src = doc.add_table(rows=12, cols=2)
    src.style = "Table Grid"
    for i, h in enumerate(["Topic", "Path in qa-automation-kb"]):
        set_cell_text(src.rows[0].cells[i], h, bold=True, color=WHITE_RGB)
    style_table_header(src.rows[0])
    for ri, (topic, path) in enumerate([
        ("Prompt H", "qa-knowledge-base/00_SYSTEM/PROMPTS.md"),
        ("Role & constraints", "qa-knowledge-base/00_SYSTEM/ROLE.md, CONSTRAINTS.md"),
        ("Glossary (triage, flakiness, RCA)", "qa-knowledge-base/00_SYSTEM/GLOSSARY.md"),
        ("Bug reporting SOP", "automation-bug-lifecycle/evidence/regression-reports/BUG_REPORTING_PROCESS.md"),
        ("Repetitive tasks guide", "qa-knowledge-base/05_ONBOARDING/HOW_TO_REPETITIVE_TASKS.md"),
        ("Defect lifecycle", "automation-bug-lifecycle/process/DEFECT_LIFECYCLE.md"),
        ("Triage rules", "automation-bug-lifecycle/process/TRIAGE_RULES.md"),
        ("Flakiness playbook", "automation-bug-lifecycle/process/FLAKINESS_PLAYBOOK.md"),
        ("RCA process", "automation-bug-lifecycle/process/RCA_PROCESS.md"),
        ("JIRA template", "automation-bug-lifecycle/templates/JIRA_TICKET_TEMPLATE.md"),
        ("Confluence exports", "automation-bug-lifecycle/reference/confluence-bug-handling/Bug Handling/"),
    ], start=1):
        set_cell_text(src.rows[ri].cells[0], topic, bold=True, color=NAVY)
        set_cell_text(src.rows[ri].cells[1], path)

    doc.save(DOCX_OUT)
    print(f"Created: {DOCX_OUT}")


# ── PPTX ──────────────────────────────────────────────────────────────────────

class Deck:
    FOOTER = "QA Automation Standard  ·  Automation Bug Lifecycle  ·  Internal"

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.page = 0

    def _slide(self):
        self.page += 1
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = PPT_OFF
        return s

    def _box(self, slide, l, t, w, h, fill, line=None, radius=True):
        st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(st, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line:
            sh.line.color.rgb = line
            sh.line.width = PptPt(1.25)
        else:
            sh.line.fill.background()
        if radius:
            sh.adjustments[0] = 0.06
        return sh

    def _txt(self, slide, l, t, w, h, text, size=16, bold=False, color=PPT_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = PptPt(4)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        return tb

    def _bullets(self, slide, items, l, t, w, h, size=15, color=PPT_TEXT, spacing=8):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"
            p.space_after = PptPt(spacing)
            p.line_spacing = 1.3
            p.level = 0
            if not item.startswith(" "):
                p.text = f"  {item}"

    def _header(self, slide, title, subtitle=""):
        self._box(slide, PptInches(0), PptInches(0), SLIDE_W, HEADER_H, PPT_NOBLE)
        self._box(slide, PptInches(0), HEADER_H, SLIDE_W, PptInches(0.07), PPT_PEAK, radius=False)
        self._txt(slide, MARGIN, PptInches(0.2), PptInches(11), PptInches(0.55),
                  title, size=30, bold=True, color=PPT_WHITE)
        if subtitle:
            self._txt(slide, MARGIN, PptInches(0.72), PptInches(11), PptInches(0.35),
                      subtitle, size=13, color=PPT_TEAL)

    def _footer(self, slide):
        self._txt(slide, MARGIN, FOOTER_Y, PptInches(10), PptInches(0.28),
                  self.FOOTER, size=9, color=PPT_MUTED)
        self._txt(slide, PptInches(12.15), FOOTER_Y, PptInches(0.7), PptInches(0.28),
                  str(self.page), size=9, color=PPT_MUTED, align=PP_ALIGN.RIGHT)

    def _fit_img(self, slide, path, l, t, max_w, max_h):
        with Image.open(path) as im:
            rw, rh = im.size
        ratio = rw / rh
        w = max_w
        h = w / ratio
        if h > max_h:
            h = max_h
            w = h * ratio
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)
        return w, h

    def _img_frame(self, slide, path, l, t, max_w, max_h, caption=""):
        pad = PptInches(0.12)
        cap_h = PptInches(0.38) if caption else PptInches(0)
        frame_h = max_h + pad * 2 + cap_h
        self._box(slide, l, t, max_w + pad * 2, frame_h, PPT_CARD, PPT_BORDER)
        iw, ih = self._fit_img(slide, path, l + pad, t + pad, max_w, max_h)
        if caption:
            self._txt(slide, l + pad, t + pad + ih + PptInches(0.06), max_w, cap_h,
                      caption, size=10, color=PPT_MUTED, align=PP_ALIGN.CENTER)

    def title(self):
        s = self._slide()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = PPT_NOBLE_DARK
        self._box(s, PptInches(0), PptInches(2.95), SLIDE_W, PptInches(0.09), PPT_PEAK, radius=False)
        self._box(s, PptInches(0), PptInches(3.04), SLIDE_W, PptInches(0.04), PPT_TEAL, radius=False)
        self._txt(s, MARGIN, PptInches(1.35), PptInches(11.5), PptInches(0.95),
                  "Automation Bug Lifecycle", size=46, bold=True, color=PPT_WHITE)
        self._txt(s, MARGIN, PptInches(2.35), PptInches(11), PptInches(0.55),
                  "QA Automation Operating Standard", size=26, color=PPT_TEAL)
        self._txt(s, MARGIN, PptInches(3.45), PptInches(9), PptInches(0.4),
                  f"Version {VERSION_DATE}  ·  Applicable to all teams", size=15, color=PPT_SLATE)
        self._txt(s, MARGIN, PptInches(4.0), PptInches(9), PptInches(0.35),
                  "qa-automation-kb  ·  AMSQUAD", size=13, color=PPT_MUTED)
        cards = [
            ("Triage", "Classify first", PPT_WARN),
            ("Evidence", "One folder", PPT_TEAL_DARK),
            ("Cursor H", "All comms", PPT_ACTION),
            ("GitLab PM", "Change set", PPT_PEAK),
        ]
        cw, gap, y = PptInches(2.65), PptInches(0.32), PptInches(5.35)
        for i, (lbl, val, col) in enumerate(cards):
            x = MARGIN + i * (cw + gap)
            self._box(s, x, y, cw, PptInches(1.2), PptRGB(0x0D, 0x2A, 0x36))
            self._box(s, x, y, cw, PptInches(0.07), col, radius=False)
            self._txt(s, x + PptInches(0.18), y + PptInches(0.2), cw, PptInches(0.3), lbl, size=10, color=PPT_MUTED)
            self._txt(s, x + PptInches(0.18), y + PptInches(0.5), cw, PptInches(0.45), val, size=19, bold=True, color=PPT_WHITE)

    def section(self, title, subtitle=""):
        s = self._slide()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = PPT_NOBLE
        self._box(s, PptInches(0), PptInches(3.5), SLIDE_W, PptInches(0.07), PPT_PEAK, radius=False)
        self._txt(s, MARGIN, PptInches(2.35), PptInches(11), PptInches(0.85),
                  title, size=38, bold=True, color=PPT_WHITE)
        if subtitle:
            self._txt(s, MARGIN, PptInches(3.72), PptInches(10), PptInches(0.5),
                      subtitle, size=18, color=PPT_TEAL)

    def content(self, title, bullets, subtitle="", size=15):
        s = self._slide()
        self._header(s, title, subtitle)
        self._bullets(s, bullets, MARGIN, BODY_TOP, SLIDE_W - MARGIN * 2, PptInches(5.4), size=size)
        self._footer(s)

    def chart(self, title, path, caption="", subtitle=""):
        s = self._slide()
        self._header(s, title, subtitle)
        self._img_frame(s, path, MARGIN, BODY_TOP, SLIDE_W - MARGIN * 2, PptInches(4.75), caption)
        self._footer(s)

    def split(self, title, bullets, path, caption="", subtitle="", img_right=True):
        s = self._slide()
        self._header(s, title, subtitle)
        tw = PptInches(5.35)
        il = PptInches(6.85) if img_right else MARGIN
        tl = MARGIN if img_right else PptInches(6.85)
        self._bullets(s, bullets, tl, BODY_TOP, tw, PptInches(5.2), size=14)
        self._img_frame(s, path, il, BODY_TOP, PptInches(5.75), PptInches(4.85), caption)
        self._footer(s)

    def cards_2x2(self, title, cards, subtitle=""):
        s = self._slide()
        self._header(s, title, subtitle)
        cw, ch, gap = PptInches(5.85), PptInches(2.35), PptInches(0.4)
        positions = [(MARGIN, BODY_TOP), (MARGIN + cw + gap, BODY_TOP),
                     (MARGIN, BODY_TOP + ch + gap), (MARGIN + cw + gap, BODY_TOP + ch + gap)]
        for (x, y), (num, heading, body, accent) in zip(positions, cards):
            self._box(s, x, y, cw, ch, PPT_CARD, PPT_BORDER)
            self._box(s, x, y, cw, PptInches(0.08), accent, radius=False)
            self._txt(s, x + PptInches(0.22), y + PptInches(0.18), PptInches(0.5), PptInches(0.4),
                      num, size=22, bold=True, color=accent)
            self._txt(s, x + PptInches(0.22), y + PptInches(0.55), cw - PptInches(0.44), PptInches(0.4),
                      heading, size=16, bold=True, color=PPT_NOBLE)
            self._txt(s, x + PptInches(0.22), y + PptInches(0.95), cw - PptInches(0.44), ch - PptInches(1.05),
                      body, size=12.5, color=PPT_SLATE)
        self._footer(s)

    def save(self, path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)


def build_pptx(charts: dict[str, Path]) -> None:
    d = Deck()
    ui = ASSETS / "gitlab_project_manager_ui.png"
    mr = ASSETS / "gitlab_mr_results_sample.png"

    d.title()

    d.content("What This Standard Defines", [
        "Repeatable response when V2 (Jenkins) or V3 (GitLab) regression fails",
        "Applies to any QA Automation team — UI, API, performance, any program",
        "Single source of truth: qa-automation-kb (00_SYSTEM through 11_BACKLOG)",
        "Integrates triage rules, Prompt H, leadership-approved comms, and GitLab change-set tooling",
        "Goal: fast, accurate footprint — not ad-hoc firefighting or noisy JIRA",
    ], subtitle="Operating norm for all teams")

    d.chart("End-to-End Workflow", charts["workflow"],
            "Detect → Triage → Evidence → Cursor H → JIRA/Notify → Change Set → Resolve",
            subtitle="Seven steps — every regression failure")

    d.chart("Step 1 — Triage Before You Log", charts["triage"],
            "Environment and flaky failures do not enter the full bug cycle",
            subtitle="Decision tree validated against TRIAGE_RULES.md + FLAKINESS_PLAYBOOK.md")

    d.section("Act on a Real Defect", "Evidence · Communications · Change set")

    d.cards_2x2("Four Pillars of the Standard", [
        ("01", "Triage & classify",
         "Per automation-bug-lifecycle/process/TRIAGE_RULES.md. Rerun locally. Env vs flaky vs defect. "
         "Do not log JIRA for noise.", PPT_WARN),
        ("02", "Evidence folder",
         "automation-bug-lifecycle/evidence/regression-reports/[MMDDYYYY]/ — screenshots, logs, test data, "
         "one bug .md per issue.", PPT_TEAL_DARK),
        ("03", "Cursor Prompt H",
         "qa-knowledge-base/00_SYSTEM/PROMPTS.md — JIRA block, failure email, Teams message, resolution "
         "placeholder in one file.", PPT_ACTION),
        ("04", "GitLab change set",
         "GitLab Project Manager — who merged, who committed, which branches between "
         "last green run and failure.", PPT_PEAK),
    ], subtitle="Validated against repository standards")

    d.chart("Evidence Folder Structure", charts["evidence"],
            "Naming: [MMDDYYYY]_[Feature]_[IssueType].md",
            subtitle="BUG_REPORTING_PROCESS.md")

    d.split("Cursor Prompt H — One Prompt, All Deliverables", [
        "Open qa-automation-kb in Cursor",
        "qa-knowledge-base/00_SYSTEM/PROMPTS.md → section H",
        "Provide: date, feature, error, report URL, folder path, file list",
        "Attach failure screenshots in chat for richer context",
        "",
        "Output sections:",
        "  • JIRA copy-paste block (Summary, Steps, Env, Priority)",
        "  • Failure email — leadership-approved To/Cc",
        "  • Teams message with JIRA + report links",
        "  • Resolution email placeholder (template 1b)",
    ], charts["toolchain"], "Prompt H toolchain", subtitle="qa-knowledge-base/05_ONBOARDING/HOW_TO_REPETITIVE_TASKS.md")

    d.content("Communication Standards", [
        "Teams and email templates are approved by leadership and senior QA resources",
        "Use Prompt H output — do not improvise distribution lists or email structure",
        "",
        "To: AGS Tech Leads, AGS Chapter Leads, AGS Development, Brian Danilczyk",
        "Cc: Rajib Akhter · Henry Dittmer · Phuong Huynh · Automation.Squad",
        "",
        "Subject: [Priority] [Feature] — [Issue Type]",
        "Sources: Confluence Bug Handling PDFs (1. Failure · 1b. Resolution)",
    ], subtitle="Leadership-approved templates")

    d.split("GitLab Project Manager — UI", [
        "Local utility: GitLabInfoProjUI (read-only GitLab API)",
        "Manage Token — PAT with read_api scope",
        "Select project + date range (last pass → failure)",
        "Tabs: Merge Requests · Commits · Branches · Pipelines",
        "Filter Merged · Search by author or branch",
        "Export CSV/Excel or Copy to clipboard",
    ], ui if ui.exists() else charts["repos"],
       "Parameter panel — project, date presets, Generate",
       subtitle="Change-set investigation tool")

    d.split("GitLab Project Manager — Sample Output", [
        "Project: monolith · Type: Merge Requests",
        "Date range: last green run → failure date",
        "Per MR: Title · Author · Reviewers · Merged by",
        "Branch flow: source → target (e.g. feature/GRF-1394 → main)",
        "Created At · Merged At (EST timestamps)",
        "Paste Change Set section into JIRA comment or failure email",
    ], mr if mr.exists() else charts["repos"],
       "Real output — 8 MRs in window, 4 merged",
       subtitle="Who changed what — application footprint", img_right=False)

    d.chart("Repository Investigation Order", charts["repos"],
            "Start monolith (app) → automation → qa-automation → prime-test-automation",
            subtitle="Cross-reference MR authors with failing test area")

    d.section("Close the Loop", "Resolution · Unlock · Knowledge capture")

    d.chart("Defect Lifecycle", charts["lifecycle"],
            "automation-bug-lifecycle/process/DEFECT_LIFECYCLE.md · TRIAGE_RULES.md",
            subtitle="JIRA states")

    d.content("Resolution Standard", [
        "Dev delivers fix or revert on affected branch",
        "QA reruns impacted scenarios (full regression optional)",
        "Update JIRA with RCA → Verified → Closed",
        "Send resolution email — same To/Cc as failure (template 1b)",
        "Unlock monolith/main + automation/main if locked",
        "Update bug .md: JIRA link, Status Closed, resolution section filled",
    ], subtitle="Confluence 1b. Automation Bug Resolution Follow-Up")

    d.content("Multi-Failure Rollup", [
        "Large nightly failures (e.g. 62/358) — group by feature/plan, not per test",
        "Build failure matrix with root-cause hints per traunch",
        "One JIRA per root cause — umbrella + children when needed",
        "Single combined email/Teams with matrix reference",
        "Example: 04202026_DailyRegression_PipelineFailureRollup.md",
    ], subtitle="When many tests fail at once")

    d.content("Repository Map — Where Standards Live", [
        "00_SYSTEM — PROMPTS.md (H), ROLE.md, CONSTRAINTS.md, GLOSSARY.md",
        "automation-bug-lifecycle/process — DEFECT_LIFECYCLE, TRIAGE_RULES, FLAKINESS_PLAYBOOK, RCA_PROCESS",
        "05_ONBOARDING — HOW_TO_REPETITIVE_TASKS.md (step-by-step)",
        "06_TEMPLATES — JIRA_TICKET_TEMPLATE, RCA_TEMPLATE",
        "automation-bug-lifecycle/evidence/regression-reports/, reference/confluence-bug-handling/",
        "This module — automation-bug-lifecycle/ (playbook + deck + generator)",
    ], subtitle="If it's not in the repo, treat as unknown — per ROLE.md")

    d.content("Key Takeaways", [
        "Triage first — not every failure is a defect",
        "One evidence folder · one Prompt H run · one bug .md per issue",
        "Communications use leadership-approved templates — no improvisation",
        "GitLab PM answers who merged and committed in the failure window",
        "Standard applies to any team, any program, going forward",
        "Full playbook: automation-bug-lifecycle/deliverables/Automation-Bug-Lifecycle-Playbook.docx",
    ], subtitle="QA Automation operating norm")

    d.save(PPTX_OUT)
    print(f"Created: {PPTX_OUT}")


def main() -> None:
    charts = generate_charts()
    for k, p in charts.items():
        print(f"Chart [{k}]: {p}")
    try:
        build_docx(charts)
    except PermissionError:
        print("Skipped DOCX — close open file and re-run.")
    build_pptx(charts)


if __name__ == "__main__":
    main()
