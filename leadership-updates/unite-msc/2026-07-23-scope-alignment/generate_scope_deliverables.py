#!/usr/bin/env python3
"""Generate Unite MSC Scope Alignment DOCX + PPTX for Jul 23 leadership meeting."""

from __future__ import annotations

from datetime import date
from pathlib import Path

import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from PIL import Image

OUT_DIR = Path(__file__).resolve().parent
ASSETS = OUT_DIR / "_assets"
DOCX_OUT = OUT_DIR / "_output" / "Unite-MSC-Scope-Alignment-2026-07-23-v2.docx"
PPTX_OUT = OUT_DIR / "_output" / "Unite-MSC-Scope-Alignment-2026-07-23-v2.pptx"
REPORT_DATE = "July 23, 2026"

# Metrics (Jul 23, 2026)
M2_IMPL, M2_TOTAL = 24, 25
M1_IMPL, M1_TOTAL = 6, 27
M2_PCT = round(100 * M2_IMPL / M2_TOTAL, 1)
M1_PCT = round(100 * M1_IMPL / M1_TOTAL, 1)

# Brand
NAVY = RGBColor(0x00, 0x30, 0x57)
TEAL = RGBColor(0x00, 0x7A, 0x8C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
GRAY = RGBColor(0x61, 0x61, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = "E8EEF4"
GREEN_BG = "E8F5E9"
AMBER_BG = "FFF3E0"
RED_BG = "FFEBEE"

PPT_NAVY = PptRGB(0x00, 0x30, 0x57)
PPT_NAVY_DARK = PptRGB(0x0A, 0x16, 0x28)
PPT_TEAL = PptRGB(0x00, 0xC2, 0xCB)
PPT_TEAL_DARK = PptRGB(0x00, 0x7A, 0x8C)
PPT_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
PPT_OFF_WHITE = PptRGB(0xF8, 0xFA, 0xFC)
PPT_SLATE = PptRGB(0x94, 0xA3, 0xB8)
PPT_TEXT = PptRGB(0x1E, 0x29, 0x3B)
PPT_TEXT_MUTED = PptRGB(0x64, 0x74, 0x8B)
PPT_GREEN = PptRGB(0x16, 0xA3, 0x4A)
PPT_GREEN_LIGHT = PptRGB(0xDC, 0xFC, 0xE7)
PPT_AMBER = PptRGB(0xD9, 0x77, 0x06)
PPT_AMBER_LIGHT = PptRGB(0xFF, 0xED, 0xD5)
PPT_CARD = PptRGB(0xF1, 0xF5, 0xF9)
PPT_CARD_BORDER = PptRGB(0xE2, 0xE8, 0xF0)

# Slide canvas (16:9 widescreen)
SLIDE_W = PptInches(13.333)
SLIDE_H = PptInches(7.5)
MARGIN = PptInches(0.55)
HEADER_H = PptInches(1.05)
CONTENT_TOP = PptInches(1.35)
FOOTER_Y = PptInches(7.05)

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Segoe UI", "Calibri", "Arial"],
})


# ── DOCX helpers ──────────────────────────────────────────────────────────────

def shade_cell(cell, hex_color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False, color: RGBColor | None = None, size: int = 10) -> None:
    cell.text = ""
    run = cell.paragraphs[0].add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color


def style_table_header(row) -> None:
    for cell in row.cells:
        shade_cell(cell, "003057")
        for p in cell.paragraphs:
            for run in p.runs:
                run.font.color.rgb = WHITE
                run.bold = True
                run.font.size = Pt(10)
                run.font.name = "Calibri"


def add_section_heading(doc: Document, title: str, level: int = 1) -> None:
    h = doc.add_heading(title, level=level)
    for run in h.runs:
        run.font.color.rgb = NAVY
        run.font.name = "Calibri"


def add_callout(doc: Document, text: str, bg: str = LIGHT_BG) -> None:
    table = doc.add_table(rows=1, cols=1)
    cell = table.rows[0].cells[0]
    shade_cell(cell, bg)
    set_cell_text(cell, text, size=11)
    for p in cell.paragraphs:
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        for run in p.runs:
            run.font.color.rgb = NAVY


def setup_header_footer(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    hp = section.header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.add_run("Unite MSC Scope Alignment  |  Internal")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"
    fp = section.footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run(f"QA Automation — AMSQUAD  |  {REPORT_DATE}")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"


# ── Charts ────────────────────────────────────────────────────────────────────

def _save(fig, name: str) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / name
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_m2_coverage(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(7, 3.8))
    verified, excluded = M2_IMPL, M2_TOTAL - M2_IMPL
    ax.barh(["Mobile 2 Business Endpoints"], [verified], color="#2E7D32", height=0.4, label=f"Implemented ({verified})")
    ax.barh(["Mobile 2 Business Endpoints"], [excluded], left=[verified], color="#B0BEC5", height=0.4, label=f"Excluded ({excluded})")
    ax.set_xlim(0, M2_TOTAL)
    ax.set_xlabel("Endpoint count (documented business scope = 25)", fontsize=10)
    ax.set_title(f"Mobile 2 API — {M2_IMPL}/{M2_TOTAL} Implemented ({M2_PCT}%)", fontsize=12, fontweight="bold", color="#003057")
    ax.axvline(M2_IMPL, color="#003057", linestyle="--", linewidth=1, alpha=0.6)
    ax.text(M2_IMPL + 0.15, 0, f"{M2_PCT}%", fontsize=11, fontweight="bold", color="#003057", va="center")
    ax.legend(loc="lower right", fontsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    note = "Excluded: GET mobilemembers/{planId}/{username} — acceptance harness only (smoke)"
    fig.text(0.5, 0.01, note, ha="center", fontsize=8, color="#616161")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_m1_coverage(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(5.5, 4))
    remaining = M1_TOTAL - M1_IMPL
    ax.pie(
        [M1_IMPL, remaining],
        explode=(0.06, 0),
        labels=[f"Implemented\n({M1_IMPL})", f"Remaining\n({remaining})"],
        colors=["#2E7D32", "#E0E0E0"],
        autopct="%1.1f%%",
        startangle=90,
        textprops={"fontsize": 10},
    )
    ax.set_title(f"Mobile 1 API — {M1_IMPL}/{M1_TOTAL} Implemented ({M1_PCT}%)\nTests on main; suite wiring ongoing",
                 fontsize=11, fontweight="bold", color="#003057")
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_validation_layers(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 5.5))
    layers = [
        ("L5 — SQL API–DB reconciliation", "Enhancement / pilot only", "#C62828", 0.55),
        ("L4 — Business data assertions", "IN SCOPE — sign-off bar", "#16A34A", 0.70),
        ("L3 — JSON schema / structure", "IN SCOPE", "#007A8C", 0.85),
        ("L2 — Response contract", "IN SCOPE", "#007A8C", 0.95),
        ("L1 — HTTP / auth / transport", "IN SCOPE", "#003057", 1.05),
    ]
    for i, (label, status, color, width) in enumerate(layers):
        ax.barh(i, width, color=color, height=0.55, alpha=0.92)
        ax.text(0.02, i, label, va="center", fontsize=11, color="white", fontweight="bold")
        ax.text(width - 0.02, i, status, va="center", ha="right", fontsize=9, color="white")
    ax.set_yticks([])
    ax.set_xlim(0, 1.15)
    ax.set_title("API Validation Layers", fontsize=15, fontweight="bold", color="#003057", pad=12)
    ax.axis("off")
    fig.patch.set_facecolor("#FAFBFC")
    fig.text(0.5, 0.02, "Mobile 2 sign-off = L1–L4. L5 is a separate program decision.",
             ha="center", fontsize=9.5, color="#64748B")
    fig.tight_layout(rect=[0.02, 0.05, 0.98, 0.98])
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="#FAFBFC")
    plt.close(fig)


def chart_hybrid_env(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(8.5, 4))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")
    boxes = [
        (0.3, 3.2, "Design &\nDevelop", "#003057"),
        (2.5, 3.2, "Execute\nTests", "#007A8C"),
        (4.7, 3.2, "Evidence &\nSign-off", "#2E7D32"),
        (7.0, 3.2, "Merge to\nmain", "#003057"),
        (1.0, 1.0, "QC4\n(unstable)", "#C62828"),
        (4.0, 1.0, "Stage 1\n(primary)", "#2E7D32"),
        (7.0, 1.0, "Perf testing\n(Stage 1 only)", "#2E7D32"),
    ]
    for x, y, text, color in boxes:
        rect = mpatches.FancyBboxPatch((x, y), 1.6, 0.9, boxstyle="round,pad=0.05",
                                       linewidth=1.5, edgecolor=color, facecolor=color, alpha=0.15)
        ax.add_patch(rect)
        ax.text(x + 0.8, y + 0.45, text, ha="center", va="center", fontsize=8.5, fontweight="bold", color=color)
    for x1, y1, x2, y2 in [(1.9, 3.65, 2.5, 3.65), (4.1, 3.65, 4.7, 3.65), (6.3, 3.65, 7.0, 3.65)]:
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1), arrowprops=dict(arrowstyle="->", color="#616161", lw=1.5))
    ax.annotate("", xy=(3.3, 3.2), xytext=(1.8, 1.9), arrowprops=dict(arrowstyle="->", color="#2E7D32", lw=1.2, linestyle="dashed"))
    ax.annotate("", xy=(5.5, 3.2), xytext=(4.8, 1.9), arrowprops=dict(arrowstyle="->", color="#2E7D32", lw=1.5))
    ax.text(5.0, 0.3, "Hybrid rule: when QC4 blocked by env/team deps, complete on Stage 1. Tests pass when auth path is healthy.",
            ha="center", fontsize=8.5, color="#616161")
    ax.set_title("Hybrid Environment Strategy — Functional + Performance", fontsize=12, fontweight="bold", color="#003057", pad=10)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_scope_options(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 5.2))
    options = ["A — L1–L4 only\n(recommended)", "B — L5 pilot\n(dashboard)", "C — Full L5\n(all endpoints)"]
    effort = [4, 8, 22]
    colors = ["#2E7D32", "#E65100", "#C62828"]
    bars = ax.bar(options, effort, color=colors, edgecolor="white", linewidth=2, width=0.52)
    ax.set_ylabel("Relative effort (weeks, illustrative)", fontsize=11, color="#334155")
    ax.set_title("SQL Validation Scope Options", fontsize=14, fontweight="bold", color="#003057", pad=14)
    ax.set_ylim(0, 26)
    ax.tick_params(colors="#475569")
    for bar, e in zip(bars, effort):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.6,
                f"~{e} wks", ha="center", fontsize=11, fontweight="bold", color="#003057")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CBD5E1")
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.set_facecolor("#FAFBFC")
    fig.patch.set_facecolor("#FAFBFC")
    fig.text(0.5, 0.02,
             "Option A = current Mobile 2 delivery. Options B/C require SME SQL mapping + multi-sprint program.",
             ha="center", fontsize=9, color="#64748B")
    fig.tight_layout(rect=[0, 0.05, 1, 1])
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="#FAFBFC")
    plt.close(fig)


def generate_charts() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    charts = {
        "m2": ASSETS / "chart_m2_coverage.png",
        "m1": ASSETS / "chart_m1_coverage.png",
        "validation": ASSETS / "chart_validation_layers.png",
        "hybrid": ASSETS / "chart_hybrid_env.png",
        "scope": ASSETS / "chart_scope_options.png",
    }
    chart_m2_coverage(charts["m2"])
    chart_m1_coverage(charts["m1"])
    chart_validation_layers(charts["validation"])
    chart_hybrid_env(charts["hybrid"])
    chart_scope_options(charts["scope"])
    return charts


# ── DOCX ──────────────────────────────────────────────────────────────────────

def build_docx(charts: dict[str, Path]) -> None:
    DOCX_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    setup_header_footer(doc)

    # Cover
    for _ in range(5):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("UNITE MSC")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = NAVY
    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = s.add_run("Automation Scope Alignment")
    r2.font.size = Pt(20)
    r2.font.color.rgb = TEAL
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for line in [REPORT_DATE, "Audience: Leadership", "Prepared by: QA Automation (AMSQUAD)",
                 "Follow-up to Jul 17 bi-weekly leadership sync"]:
        run = meta.add_run(line + "\n")
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY
    doc.add_page_break()

    # TOC
    add_section_heading(doc, "Table of Contents", level=1)
    for item in [
        "1. Meeting Purpose",
        "2. Program Status",
        "3. API Validation — Scope Clarification",
        "4. SQL Validation — Discussion Topic",
        "5. Hybrid Environment Strategy",
        "6. Mobile 2 — Delivery Summary",
        "7. Mobile 1 — Delivered So Far",
        "8. Performance Testing — Stage 1 Approach",
        "9. Decisions & Next Steps",
        "Appendix A — SQL Scope Options (discussion reference)",
        "Appendix B — Jul 17 Open Questions",
    ]:
        doc.add_paragraph(item, style="List Number")
    doc.add_page_break()

    # 1 Purpose
    add_section_heading(doc, "1. Meeting Purpose", level=1)
    add_callout(doc,
        "Align leadership on Unite MSC automation scope — especially SQL validation depth, "
        "environment strategy (QC4 vs Stage 1), and what constitutes 'complete' for Mobile 2 and Mobile 1.")
    doc.add_paragraph(
        "The Jul 17 bi-weekly raised that SQL API–DB validation could significantly expand scope. "
        "This document separates what is delivered today (L1–L4) from optional enhancement work (L5), "
        "and documents the hybrid environment approach while QC4 remains unstable.")
    add_section_heading(doc, "Decisions needed", level=2)
    for item in [
        "Confirm L1–L4 as the Mobile 2 sign-off bar (no L5 required for 'complete')",
        "If L5 is desired — choose pilot (dashboard) vs full program vs defer",
        "Confirm hybrid Stage 1 strategy for functional + performance evidence",
        "Identify SMEs for sign-off and production-like performance criteria",
        "Confirm strategic direction and reusability expectations for MSC API automation",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # 2 Status
    add_section_heading(doc, "2. Program Status — Updated Jul 23", level=1)
    add_callout(doc, f"Mobile 2 is COMPLETE at {M2_IMPL}/{M2_TOTAL} ({M2_PCT}%) business endpoints. "
                 f"Mobile 1 is IN PROGRESS at {M1_IMPL}/{M1_TOTAL} ({M1_PCT}%).", GREEN_BG)

    kpi_table = doc.add_table(rows=2, cols=4)
    kpi_table.style = "Table Grid"
    kpi_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(["Mobile 2", "Mobile 1", "Sign-off bar", "Primary environment"]):
        set_cell_text(kpi_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(kpi_table.rows[0])
    kpi_vals = [
        (f"{M2_IMPL}/{M2_TOTAL} ({M2_PCT}%)", "COMPLETE"),
        (f"{M1_IMPL}/{M1_TOTAL} ({M1_PCT}%)", "IN PROGRESS"),
        ("L1–L4", "Delivered"),
        ("Stage 1", "Functional + perf"),
    ]
    for ci, (val, _) in enumerate(kpi_vals):
        shade_cell(kpi_table.rows[1].cells[ci], GREEN_BG if ci == 0 else (AMBER_BG if ci == 1 else LIGHT_BG))
        set_cell_text(kpi_table.rows[1].cells[ci], val, bold=True, color=NAVY)
    doc.add_paragraph()

    doc.add_picture(str(charts["m2"]), width=Inches(6.2))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_picture(str(charts["m1"]), width=Inches(4.8))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    status_table = doc.add_table(rows=5, cols=4)
    status_table.style = "Table Grid"
    headers = ["Track", "Implemented", "Scope", "Status"]
    for i, h in enumerate(headers):
        set_cell_text(status_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(status_table.rows[0])
    rows = [
        ("Mobile 2 API", f"{M2_IMPL}/{M2_TOTAL}", f"{M2_PCT}%", "COMPLETE — sign-off path"),
        ("Mobile 1 API", f"{M1_IMPL}/{M1_TOTAL}", f"{M1_PCT}%", "IN PROGRESS — tests on main"),
        ("Performance (MSC)", "1 flow live", "Stage 1", "Expanding — Stage 1 only"),
        ("Pipeline (GHA)", "Dashboard slice", "—", "Module suites in progress"),
    ]
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            shade_cell(status_table.rows[ri].cells[ci], LIGHT_BG if ci == 0 else "FFFFFF")
            set_cell_text(status_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)

    p = doc.add_paragraph()
    p.add_run("Counting rule: ").bold = True
    p.add_run("Endpoint is 'implemented' when canonical test class exists on main — even if not yet wired "
              "to master regression suite (weekly hygiene adds suite links).")
    doc.add_page_break()

    # 3 Validation
    add_section_heading(doc, "3. API Validation — Scope Clarification", level=1)
    doc.add_picture(str(charts["validation"]), width=Inches(6.2))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    val_table = doc.add_table(rows=6, cols=4)
    val_table.style = "Table Grid"
    for i, h in enumerate(["Layer", "What it validates", "Mobile 2 status", "Sign-off?"]):
        set_cell_text(val_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(val_table.rows[0])
    val_rows = [
        ("L1 HTTP / Auth", "Status codes, session/token, transport errors", "Delivered", "Yes"),
        ("L2 Response contract", "Content-type, JSON parse, error payloads", "Delivered", "Yes"),
        ("L3 Schema / structure", "JSON schema, required fields, types", "Delivered", "Yes"),
        ("L4 Business assertions", "Field values, business rules, branding", "Delivered", "Yes"),
        ("L5 SQL API–DB", "JDBC compare to API JSON via field maps", "Not delivered", "Decision needed"),
    ]
    bg_map = {"Yes": GREEN_BG, "Decision needed": AMBER_BG}
    for ri, row in enumerate(val_rows, start=1):
        for ci, val in enumerate(row):
            bg = bg_map.get(val, LIGHT_BG if ci == 0 else "FFFFFF")
            shade_cell(val_table.rows[ri].cells[ci], bg)
            set_cell_text(val_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)

    add_section_heading(doc, "Why L5 SQL is complex for Unite MSC", level=2)
    for item in [
        "Mobile APIs are BFF-layer — data is assembled from multiple SQL sources, not 1:1 endpoint-to-query",
        "Application team also confirmed indirect mapping is difficult",
        "Computed fields (totalBalance, displayInStackup, regType display names) require BFF logic replication",
        "Some fields are on-prem only or external (CMS content) — no SQL compare possible",
        "L5 is valuable as a pilot/enhancement — but it is a separate multi-sprint program, not a sign-off blocker",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # 4 SQL discussion
    add_section_heading(doc, "4. SQL Validation — Discussion Topic", level=1)
    add_callout(doc,
        "Scope options chart is retained as Appendix A for verbal discussion — it is not included in the presentation deck.",
        AMBER_BG)
    for item in [
        "Current delivery: HTTP + schema + business assertions (L1–L4)",
        "L5 SQL API–DB compare = JDBC queries mapped to API JSON fields",
        "Unite MSC BFF layer makes 1:1 SQL mapping very difficult",
        "Development team also struggles with indirect SQL-to-endpoint mapping",
        "Recommendation: L5 as optional pilot — NOT a Mobile 2 blocker",
        "Leadership to decide: defer L5, pilot on dashboard only, or approve full L5 program",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # 5 Hybrid env
    add_section_heading(doc, "5. Hybrid Environment Strategy", level=1)
    add_callout(doc,
        "QC4 is currently unstable due to team dependencies (OKD disabled during Stage 1 migration). "
        "We use Stage 1 as the primary execution environment for design, test, sign-off, and performance.",
        AMBER_BG)
    doc.add_picture(str(charts["hybrid"]), width=Inches(6.5))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    env_table = doc.add_table(rows=5, cols=4)
    env_table.style = "Table Grid"
    for i, h in enumerate(["Environment", "Functional tests", "Performance", "Notes"]):
        set_cell_text(env_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(env_table.rows[0])
    env_rows = [
        ("QC4", "When stable", "Not primary", "Some OKD failures env-specific; NMD paths work"),
        ("Stage 1", "Primary", "Primary", "Verified Jul 22 — tests pass when auth healthy"),
        ("Stage 5 / 2", "Future", "TBD", "Framework supports via Maven profiles"),
        ("Production-like load", "N/A", "Stage 1", "25-user baseline; SME criteria pending"),
    ]
    for ri, row in enumerate(env_rows, start=1):
        for ci, val in enumerate(row):
            shade_cell(env_table.rows[ri].cells[ci], LIGHT_BG if ci == 0 else "FFFFFF")
            set_cell_text(env_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)
    doc.add_page_break()

    # 6 Mobile 2
    add_section_heading(doc, "6. Mobile 2 — Delivery Summary", level=1)
    add_callout(doc, f"All {M2_IMPL} business endpoints have canonical TestNG tests on main. "
                 f"1 harness endpoint excluded by design.", GREEN_BG)
    m2_detail = [
        ("Framework", "Canonical TestNG — shared auth, HTML reports, OKD/NMD branding, module + master suites"),
        ("Dashboard", "mobiledashboard + mobileytdsummary — complete"),
        ("Banks", "List, GET-by-id, POST on master; PUT/DELETE smoke-only (destructive)"),
        ("Contribution", "6 endpoints; DELETE module-only; detail/PUT Stage 1 fixture known"),
        ("Content / Plans / Activity / Investment / Balance / UGift", "Complete on master regression"),
        ("Excluded", "GET mobilemembers/{planId}/{username} — acceptance harness (smoke only)"),
        ("Evidence", "Stage 1 execution verified Jul 22; QC4 re-run when env stable"),
    ]
    for label, detail in m2_detail:
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(detail)

    # 7 Mobile 1
    add_section_heading(doc, "7. Mobile 1 — Delivered So Far", level=1)
    m1_table = doc.add_table(rows=7, cols=3)
    m1_table.style = "Table Grid"
    for i, h in enumerate(["Endpoint", "Test class", "Status"]):
        set_cell_text(m1_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(m1_table.rows[0])
    m1_rows = [
        ("POST mobilemembersession", "Mobile1AuthenticationTest", "Complete — auth foundation"),
        ("GET mobileowner", "MobileOwnerRequestTest", "Implemented on main"),
        ("GET mobileOwnerMenu", "MobileOwnerMenuRequestTest", "Implemented on main"),
        ("GET mobileprofilemenu", "MobileProfileMenuRequestTest", "Implemented on main"),
        ("GET mobilebeneficiaryByExt/{ext}", "MobileBeneficiaryByExtRequestTest", "Implemented on main"),
        ("GET mobilebankinfobyroutingnum/{routingNum}", "MobileBankInfoByRoutingNumRequestTest", "Implemented on main"),
    ]
    for ri, row in enumerate(m1_rows, start=1):
        for ci, val in enumerate(row):
            shade_cell(m1_table.rows[ri].cells[ci], GREEN_BG if ci == 2 else (LIGHT_BG if ci == 0 else "FFFFFF"))
            set_cell_text(m1_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)
    doc.add_paragraph(f"Remaining: {M1_TOTAL - M1_IMPL} endpoints in active migration.")

    # 8 Performance
    add_section_heading(doc, "8. Performance Testing — Stage 1 Approach", level=1)
    for item in [
        "Stage 1 only — stable, production-like load possible",
        "Regulations + DoD established from scratch; nightly Jenkins regression live (weekdays)",
        "Baseline: 25 users; BlazeMeter reports",
        "SME input needed for production-like success criteria",
        "Expansion path: Dashboard → Contribution → Banks → Activity flows",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_page_break()

    # 9 Decisions
    add_section_heading(doc, "9. Decisions & Next Steps", level=1)
    dec_table = doc.add_table(rows=6, cols=3)
    dec_table.style = "Table Grid"
    for i, h in enumerate(["#", "Decision / Action", "Owner"]):
        set_cell_text(dec_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(dec_table.rows[0])
    decisions = [
        ("1", "Confirm L1–L4 as Mobile 2 sign-off bar (no L5 required)", "Leadership"),
        ("2", "If L5 desired — approve pilot scope (dashboard only) or defer", "Leadership"),
        ("3", "Approve hybrid Stage 1 strategy while QC4 unstable", "Leadership"),
        ("4", "Identify SMEs for Mobile 2 sign-off + perf criteria", "Program SMEs"),
        ("5", "Confirm MSC automation strategic direction & reusability", "Leadership"),
    ]
    for ri, row in enumerate(decisions, start=1):
        for ci, val in enumerate(row):
            shade_cell(dec_table.rows[ri].cells[ci], LIGHT_BG if ci == 0 else "FFFFFF")
            set_cell_text(dec_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)

    doc.add_page_break()

    # Appendix A
    add_section_heading(doc, "Appendix A — SQL Scope Options (Discussion Reference)", level=1)
    add_callout(doc,
        "This chart supports verbal discussion only. It is not included in the presentation deck.",
        AMBER_BG)
    doc.add_picture(str(charts["scope"]), width=Inches(6.0))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    scope_table = doc.add_table(rows=4, cols=3)
    scope_table.style = "Table Grid"
    for i, h in enumerate(["Option", "Scope", "Relative effort"]):
        set_cell_text(scope_table.rows[0].cells[i], h, bold=True, color=WHITE)
    style_table_header(scope_table.rows[0])
    for ri, row in enumerate([
        ("A (recommended)", "L1–L4 only — current Mobile 2 delivery", "~4 weeks"),
        ("B", "L5 pilot — dashboard endpoint only", "~8 weeks"),
        ("C", "Full L5 — all endpoints", "~22 weeks"),
    ], start=1):
        for ci, val in enumerate(row):
            shade_cell(scope_table.rows[ri].cells[ci], GREEN_BG if ri == 1 else "FFFFFF")
            set_cell_text(scope_table.rows[ri].cells[ci], val, bold=(ci == 0), color=NAVY if ci == 0 else GRAY)

    add_section_heading(doc, "Appendix B — Jul 17 Open Questions Carried Forward", level=1)
    for q in [
        "Value and reusability of United MSC API automation",
        "SQL validation depth and approach",
        "SME input for performance metrics and success criteria",
        "SME identification for sign-off",
        "Framework focus and strategic direction",
    ]:
        doc.add_paragraph(q, style="List Bullet")

    doc.save(DOCX_OUT)
    print(f"Created: {DOCX_OUT}")


# ── PPTX — modern widescreen deck ─────────────────────────────────────────────

class SlideDesigner:
    """Gamma-style executive deck: navy header, light body, card layouts, safe image fitting."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._page = 0

    def _blank(self):
        self._page += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PPT_OFF_WHITE
        return slide

    def _rect(self, slide, left, top, width, height, fill, line=None, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line:
            shape.line.color.rgb = line
            shape.line.width = PptPt(1)
        else:
            shape.line.fill.background()
        if radius:
            shape.adjustments[0] = 0.08
        return shape

    def _text(self, slide, left, top, width, height, text, size=18, bold=False,
              color=PPT_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.alignment = align
        return box

    def _header(self, slide, title: str, subtitle: str = "") -> None:
        self._rect(slide, PptInches(0), PptInches(0), SLIDE_W, HEADER_H, PPT_NAVY)
        self._rect(slide, PptInches(0), HEADER_H, SLIDE_W, PptInches(0.06), PPT_TEAL)
        self._text(slide, MARGIN, PptInches(0.22), PptInches(10), PptInches(0.55),
                   title, size=28, bold=True, color=PPT_WHITE)
        if subtitle:
            self._text(slide, MARGIN, PptInches(0.68), PptInches(10), PptInches(0.35),
                       subtitle, size=13, color=PPT_TEAL)

    def _footer(self, slide) -> None:
        self._text(slide, MARGIN, FOOTER_Y, PptInches(6), PptInches(0.3),
                   "Unite MSC Automation  ·  Scope Alignment  ·  Internal", size=9, color=PPT_TEXT_MUTED)
        self._text(slide, PptInches(12.2), FOOTER_Y, PptInches(0.8), PptInches(0.3),
                   str(self._page), size=9, color=PPT_TEXT_MUTED, align=PP_ALIGN.RIGHT)

    def _bullets(self, slide, items: list[str], left, top, width, height,
                 size=15, color=PPT_TEXT, spacing=6) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"●  {item}"
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"
            p.space_after = PptPt(spacing)
            p.line_spacing = 1.25

    def _fitted_image(self, slide, path: Path, left, top, max_w, max_h):
        with Image.open(path) as img:
            w, h = img.size
        ratio = w / h
        width = max_w
        height = width / ratio
        if height > max_h:
            height = max_h
            width = height * ratio
        pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return pic, width, height

    def _image_card(self, slide, path: Path, left, top, max_w, max_h, caption: str = "") -> None:
        pad = PptInches(0.18)
        card = self._rect(slide, left, top, max_w + pad * 2, max_h + pad * 2 + (PptInches(0.35) if caption else 0),
                          PPT_WHITE, PPT_CARD_BORDER, radius=True)
        pic, w, h = self._fitted_image(slide, path, left + pad, top + pad, max_w, max_h)
        if caption:
            self._text(slide, left + pad, top + pad + h + PptInches(0.08), max_w,
                       PptInches(0.28), caption, size=10, color=PPT_TEXT_MUTED, align=PP_ALIGN.CENTER)

    def _metric_card(self, slide, left, top, width, height, label, value, accent=PPT_TEAL_DARK) -> None:
        self._rect(slide, left, top, width, height, PPT_WHITE, PPT_CARD_BORDER, radius=True)
        self._rect(slide, left, top, width, PptInches(0.07), accent)
        self._text(slide, left + PptInches(0.2), top + PptInches(0.22), width - PptInches(0.4), PptInches(0.35),
                   label, size=11, color=PPT_TEXT_MUTED)
        self._text(slide, left + PptInches(0.2), top + PptInches(0.55), width - PptInches(0.4), PptInches(0.55),
                   value, size=26, bold=True, color=PPT_NAVY)

    def title_slide(self) -> None:
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PPT_NAVY_DARK
        # Accent gradient bar
        self._rect(slide, PptInches(0), PptInches(3.1), SLIDE_W, PptInches(0.08), PPT_TEAL)
        self._text(slide, MARGIN, PptInches(1.6), PptInches(11), PptInches(0.9),
                   "Unite MSC Automation", size=44, bold=True, color=PPT_WHITE)
        self._text(slide, MARGIN, PptInches(2.55), PptInches(11), PptInches(0.55),
                   "Scope Alignment", size=28, color=PPT_TEAL)
        self._text(slide, MARGIN, PptInches(3.55), PptInches(8), PptInches(0.4),
                   REPORT_DATE, size=16, color=PPT_SLATE)
        self._text(slide, MARGIN, PptInches(4.1), PptInches(8), PptInches(0.35),
                   "QA Automation  ·  AMSQUAD", size=14, color=PPT_SLATE)
        # KPI strip at bottom
        card_w = PptInches(2.7)
        gap = PptInches(0.35)
        start_x = MARGIN
        y = PptInches(5.5)
        for i, (lbl, val, col) in enumerate([
            ("Mobile 2", f"{M2_IMPL}/{M2_TOTAL}", PPT_GREEN),
            ("Mobile 1", f"{M1_IMPL}/{M1_TOTAL}", PPT_AMBER),
            ("Sign-off", "L1–L4", PPT_TEAL_DARK),
            ("Environment", "Stage 1", PPT_TEAL_DARK),
        ]):
            x = start_x + i * (card_w + gap)
            self._rect(slide, x, y, card_w, PptInches(1.15), PptRGB(0x1E, 0x3A, 0x52), radius=True)
            self._rect(slide, x, y, card_w, PptInches(0.06), col)
            self._text(slide, x + PptInches(0.15), y + PptInches(0.18), card_w, PptInches(0.3),
                       lbl, size=10, color=PPT_SLATE)
            self._text(slide, x + PptInches(0.15), y + PptInches(0.48), card_w, PptInches(0.45),
                       val, size=22, bold=True, color=PPT_WHITE)

    def section_slide(self, title: str, subtitle: str = "") -> None:
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PPT_NAVY
        self._rect(slide, PptInches(0), PptInches(3.55), SLIDE_W, PptInches(0.06), PPT_TEAL)
        self._text(slide, MARGIN, PptInches(2.4), PptInches(11), PptInches(0.8),
                   title, size=36, bold=True, color=PPT_WHITE)
        if subtitle:
            self._text(slide, MARGIN, PptInches(3.75), PptInches(10), PptInches(0.5),
                       subtitle, size=18, color=PPT_TEAL)

    def content_slide(self, title: str, bullets: list[str], subtitle: str = "") -> None:
        slide = self._blank()
        self._header(slide, title, subtitle)
        self._bullets(slide, bullets, MARGIN, CONTENT_TOP, SLIDE_W - MARGIN * 2,
                      PptInches(5.5), size=16)
        self._footer(slide)

    def split_slide(self, title: str, bullets: list[str], image: Path,
                    caption: str = "", img_side: str = "right") -> None:
        slide = self._blank()
        self._header(slide, title)
        text_w = PptInches(5.6)
        img_max_w = PptInches(5.8)
        img_max_h = PptInches(4.8)
        if img_side == "right":
            text_left = MARGIN
            img_left = PptInches(6.9)
        else:
            text_left = PptInches(6.9)
            img_left = MARGIN
        self._bullets(slide, bullets, text_left, CONTENT_TOP, text_w, PptInches(5.2), size=15)
        self._image_card(slide, image, img_left, CONTENT_TOP, img_max_w, img_max_h, caption)
        self._footer(slide)

    def chart_slide(self, title: str, image: Path, caption: str = "", subtitle: str = "") -> None:
        slide = self._blank()
        self._header(slide, title, subtitle)
        img_left = MARGIN
        img_max_w = SLIDE_W - MARGIN * 2
        img_max_h = PptInches(4.85)
        self._image_card(slide, image, img_left, CONTENT_TOP, img_max_w, img_max_h, caption)
        self._footer(slide)

    def metrics_slide(self, title: str, charts: dict) -> None:
        slide = self._blank()
        self._header(slide, title, "Updated July 23, 2026")
        cw = PptInches(2.85)
        ch = PptInches(1.35)
        gap = PptInches(0.3)
        cards = [
            ("Mobile 2 API", f"{M2_IMPL}/{M2_TOTAL}", f"{M2_PCT}%", "COMPLETE", PPT_GREEN),
            ("Mobile 1 API", f"{M1_IMPL}/{M1_TOTAL}", f"{M1_PCT}%", "IN PROGRESS", PPT_AMBER),
            ("Performance", "Live", "Stage 1", "Nightly Jenkins", PPT_TEAL_DARK),
            ("Pipeline", "GHA", "Dashboard", "Slice validated", PPT_TEAL_DARK),
        ]
        for i, (lbl, v1, v2, status, col) in enumerate(cards):
            x = MARGIN + i * (cw + gap)
            y = CONTENT_TOP
            self._metric_card(slide, x, y, cw, ch, lbl, v1, col)
            self._text(slide, x + PptInches(0.2), y + PptInches(0.95), cw, PptInches(0.25),
                       f"{v2}  ·  {status}", size=10, color=PPT_TEXT_MUTED)
        self._image_card(slide, charts["m2"], MARGIN, PptInches(2.95),
                         SLIDE_W - MARGIN * 2, PptInches(3.5),
                         "24 business endpoints implemented — 1 harness endpoint excluded by design")
        self._footer(slide)

    def decisions_slide(self) -> None:
        slide = self._blank()
        self._header(slide, "Decisions Needed Today", "Scope · Environment · SMEs")
        items = [
            ("01", "Confirm L1–L4 as Mobile 2 sign-off bar", PPT_GREEN_LIGHT, PPT_GREEN),
            ("02", "L5 SQL — defer, pilot (dashboard), or full program?", PPT_AMBER_LIGHT, PPT_AMBER),
            ("03", "Approve hybrid Stage 1 strategy while QC4 unstable", PPT_CARD, PPT_TEAL_DARK),
            ("04", "Identify sign-off + performance SMEs", PPT_CARD, PPT_TEAL_DARK),
            ("05", "Confirm MSC automation strategic direction", PPT_CARD, PPT_TEAL_DARK),
        ]
        cw = PptInches(3.75)
        ch = PptInches(1.05)
        gap_x = PptInches(0.35)
        gap_y = PptInches(0.3)
        for i, (num, text, bg, accent) in enumerate(items):
            col = i % 2
            row = i // 2
            x = MARGIN + col * (cw + gap_x)
            y = CONTENT_TOP + row * (ch + gap_y)
            if i == 4:
                x = MARGIN + (SLIDE_W - MARGIN * 2 - cw) / 2
                y = CONTENT_TOP + 2 * (ch + gap_y)
            self._rect(slide, x, y, cw, ch, bg, PPT_CARD_BORDER, radius=True)
            self._rect(slide, x, y, PptInches(0.55), ch, accent)
            self._text(slide, x + PptInches(0.12), y + PptInches(0.12), PptInches(0.4), PptInches(0.4),
                       num, size=14, bold=True, color=PPT_WHITE if bg != PPT_CARD else accent,
                       align=PP_ALIGN.CENTER)
            self._text(slide, x + PptInches(0.65), y + PptInches(0.22), cw - PptInches(0.8), ch - PptInches(0.3),
                       text, size=13, color=PPT_TEXT)
        self._footer(slide)

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)


def build_pptx(charts: dict[str, Path]) -> None:
    d = SlideDesigner()

    d.title_slide()

    d.content_slide("Why We Are Here", [
        "Follow-up from Jul 17 bi-weekly — SQL validation scope raised",
        "Leadership asked: how deep should Unite MSC automation go?",
        "Align on what is DELIVERED vs OPTIONAL enhancement",
        "Confirm hybrid environment strategy (QC4 unstable → Stage 1 primary)",
        "Identify SMEs for sign-off and performance criteria",
    ], subtitle="Meeting objective")

    d.metrics_slide("Program Status", charts)

    d.split_slide("Mobile 1 Progress", [
        f"{M1_IMPL} of {M1_TOTAL} endpoints implemented ({M1_PCT}%)",
        "Auth foundation complete on main",
        "Owner, profile, beneficiary, bank info delivered",
        "Tests exist — suite wiring is weekly hygiene",
        "21 endpoints remaining in active migration",
    ], charts["m1"], "Implemented = canonical test on main", img_side="right")

    d.section_slide("Validation Scope", "What we deliver vs. optional enhancement")

    d.chart_slide(
        "API Validation Layers",
        charts["validation"],
        "Sign-off bar = L1–L4  ·  L5 SQL is a separate program decision",
        subtitle="Addressing the SQL validation scope concern",
    )

    d.content_slide("SQL Validation — Discussion Topic", [
        "Current delivery: HTTP + schema + business assertions (L1–L4)",
        "L5 SQL API–DB compare = JDBC queries mapped to API JSON fields",
        "Unite MSC BFF layer makes 1:1 SQL mapping very difficult",
        "Development team also struggles with indirect SQL-to-endpoint mapping",
        "Recommendation: L5 as optional pilot — NOT a Mobile 2 blocker",
        "Scope options chart available as reference asset (not in deck)",
    ], subtitle="For leadership discussion")

    d.chart_slide(
        "Hybrid Environment Strategy",
        charts["hybrid"],
        "Stage 1 primary for functional + performance while QC4 has team dependencies",
        subtitle="QC4 instability workaround",
    )

    d.split_slide("Mobile 2 — Delivered", [
        "24/25 business endpoints complete",
        "Framework + dashboard + custom reporting",
        "OKD + NMD branding; module + master suites",
        "Excluded: mobilemembers harness (smoke only)",
        "Stage 1 verified Jul 22 when auth path healthy",
    ], charts["m2"], img_side="left")

    d.content_slide("Performance — Stage 1 Only", [
        "Stage 1 only — stable, production-like load possible",
        "Regulations + DoD established; nightly Jenkins live (weekdays)",
        "25-user baseline; BlazeMeter reports",
        "SME input needed for production-like success criteria",
        "Expansion: Dashboard → Contribution → Banks → Activity",
    ], subtitle="Performance regression program")

    d.decisions_slide()

    d.save(PPTX_OUT)
    print(f"Created: {PPTX_OUT}")


def main() -> None:
    charts = generate_charts()
    for k, p in charts.items():
        print(f"Chart [{k}]: {p}")
    try:
        build_docx(charts)
    except PermissionError:
        print("Skipped DOCX — file may be open. Close it and re-run to regenerate.")
    build_pptx(charts)


if __name__ == "__main__":
    main()
